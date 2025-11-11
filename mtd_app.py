import streamlit as st
st.set_page_config(page_title="MTD Framework", page_icon="📊", layout="wide")

import pandas as pd
import numpy as np
from queryhelper import *
import sys
from concurrent.futures import ThreadPoolExecutor,as_completed
import datetime
import copy
import json
import warnings
from datetime import timedelta
from dateutil.relativedelta import relativedelta
from sklearn.preprocessing import KBinsDiscretizer
import math
import io
import os
import re
import sys
import zipfile
import shutil
import tempfile
import datetime as dt
from pathlib import Path
from typing import Dict, List



from pptx import Presentation
warnings.filterwarnings('ignore')

# Note: Removed import-time database query to improve startup performance
# If needed, load configuration data lazily when required

def extract_build_flags(cfg_obj,
                        default_platforms=('blinkit','instamart','zepto')):
    if isinstance(cfg_obj, str):
        try:
            cfg = json.loads(cfg_obj)
        except json.JSONDecodeError:
            cfg = {}
    elif isinstance(cfg_obj, dict):
        cfg = cfg_obj
    else:
        cfg = {}

    vars_list = cfg.get('dbt', {}).get('variables', [])
    config_name = cfg.get('redshift','').get('database','')
    
    flags = {}
    for v in vars_list:
        name = v.get('name','')
        if name.startswith('build_normalized_offtake_'):
            platform = name.split('build_normalized_offtake_')[-1]
            flags[platform] = bool(v.get('value', False))


    if not flags:
        flags = {p: False for p in default_platforms}
    else:
        for p in default_platforms:
            flags.setdefault(p, False)

    return {
        'name': config_name,
        **{f'build_{p}': flags[p] for p in default_platforms}
    }
# flags_df = df['job_config'] \
#     .apply(extract_build_flags) \
#     .apply(pd.Series)

# df = pd.concat([df, flags_df], axis=1)
# df=df[['name','build_blinkit','build_instamart','build_zepto']]
# acc_map=read_google_sheet("https://docs.google.com/spreadsheets/d/12bF5-Grg4QYgIC5cRDcVONBV5lXV0-KDH1s7VN-z5F4/edit?gid=1758069998#gid=1758069998","ACC")
# acc_map=acc_map[acc_map['name'].notna()]
# acc_map=acc_map[acc_map['name']!="#N/A"].reset_index(drop=True)
# acc_map=acc_map[['workspace_id','account_id','db_name','name','chat_id']]
# df[df['name'].str.startswith('pee')]
# acc_map=acc_map.merge(df,how='left',left_on='db_name',right_on='name').fillna(False).drop(columns=['name_y']).rename(columns={'name_x':'name'})


def fetch_offtake_uploaded_till(acc):
    db_name = acc['db_name']
    return sfFetch(
        f"""SELECT lower(gc_platform) as gc_platform,max(platform_offtake_uploaded_till) as last_updated
        FROM {db_name}.aggregate.qcom_pid_metrics group by 1"""
    )
def fetch_for_acc(acc, l1_start, l1_end, l2_start, l2_end):

  """
   Fetch off-take and market share data for a given account over two specified date ranges.

    This function queries the Snowflake database for off-take metrics (`estimated_offtake_mrp` and
    `reported_offtake_mrp` or `normalized_offtake` for Blinkit accounts) across two periods
    ('dr1' and 'dr2'), aggregates the results at the category and brand level, and computes
    market share percentages. It supports a special Blinkit mode controlled by the
    `build_blinkit` flag in the account configuration.

    Parameters:
        acc (dict): Account configuration with keys:
            - 'db_name' (str): Name of the Snowflake database/schema to query.
            - 'build_blinkit' (bool, optional): If True, use Blinkit-specific off-take logic.
        l1_start (datetime): Start date of the first period.
        l1_end (datetime): End date of the first period.
        l2_start (datetime): Start date of the second period.
        l2_end (datetime): End date of the second period.

    Returns:
        pandas.DataFrame: A DataFrame with columns:
            - date_period (str): 'dr1' or 'dr2' indicating the period.
            - gc_platform (str): Lowercased platform identifier.
            - bgr (str): Lowercased brand group.
            - gc_city (str): City of record.
            - brand (str): Brand name.
            - actual_sum (float): Sum of off-take values for the brand.
            - category_size (float): Total off-take in the category.
            - market_share (float): Percentage share of the category.
            - account (str): Database name used for the query."""
  
  db_name = acc['db_name']
  build_blinkit = bool(acc.get('build_blinkit', False))
  try:
      if not build_blinkit:
          query = f"""
          WITH filtered AS (
            SELECT
              CASE
                WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                      AND DATE '{l1_end.date()}'
                  THEN 'dr1'
                WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                      AND DATE '{l2_end.date()}'
                  THEN 'dr2'
                ELSE NULL
              END AS date_period,
              lower(qpm.gc_platform) as gc_platform,
              lower(qpd.bgr)              AS bgr,
              qpm.gc_city,
              qpm.brand,
              qpm.reported_offtake_mrp,
              qpm.estimated_offtake_mrp,
              qpm.is_own_brand
            FROM {db_name}.AGGREGATE.qcom_pid_metrics qpm
            LEFT JOIN {db_name}.AGGREGATE.qcom_product_dimensions qpd
              ON qpm.product_id = qpd.product_id
            WHERE 
    (
    qpm.snapshot_date 
    BETWEEN DATE '{l2_start.date()}' 
        AND LEAST(
              dateadd(month,-1,qpm.platform_offtake_uploaded_till), 
              DATE '{l2_end.date()}'
            )
  )
  OR
  (
    qpm.snapshot_date 
    BETWEEN DATE '{l1_start.date()}' 
        AND LEAST(
              qpm.platform_offtake_uploaded_till, 
              DATE '{l1_end.date()}'
            )
  )
          SELECT
            date_period,
            gc_platform,
            bgr,
            gc_city,
            brand,
            SUM(estimated_offtake_mrp) AS actual_sum,
            SUM(reported_offtake_mrp)  AS actualsale,
            SUM(SUM(estimated_offtake_mrp)) 
              OVER (PARTITION BY date_period, gc_platform, bgr, gc_city)
              AS category_size
          FROM filtered
          WHERE date_period IS NOT NULL
          GROUP BY date_period, gc_platform, bgr,gc_city,brand
          """
          df = sfFetch(query)
          df.columns = df.columns.str.lower()
          
          # compute market_share & adjust category_size back
          df['market_share'] = 100 * df['actual_sum'] / df['category_size']
          # category_size was derived from reported_offtake; revert if you need the same formula as before
          df['category_size'] = df['actualsale'] / (df['market_share']/100)

          result = df.drop(columns='actualsale')
          result['account'] = db_name

      else:
          query2 = f"""
          WITH filtered AS (
            SELECT
              CASE
                WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                      AND DATE '{l1_end.date()}'
                  THEN 'dr1'
                WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                      AND DATE '{l2_end.date()}'
                  THEN 'dr2'
                ELSE NULL
              END AS date_period,
              lower(qpm.gc_platform) as gc_platform,
              lower(qpd.bgr)         as bgr,
              qpm.gc_city,
              qpm.brand,
              qpm.is_own_brand,
              qpm.normalized_offtake,
              qpm.reported_offtake_mrp
            FROM {db_name}.AGGREGATE.qcom_pid_metrics qpm
            LEFT JOIN {db_name}.AGGREGATE.qcom_product_dimensions qpd
              ON qpm.product_id = qpd.product_id
            WHERE (
    qpm.snapshot_date 
    BETWEEN DATE '{l2_start.date()}' 
        AND LEAST(
              dateadd(month,-1,qpm.platform_offtake_uploaded_till), 
              DATE '{l2_end.date()}'
            )
  )
  OR
  (
    qpm.snapshot_date 
    BETWEEN DATE '{l1_start.date()}' 
        AND LEAST(
              qpm.platform_offtake_uploaded_till, 
              DATE '{l1_end.date()}'
            )
          )),
          agg as(SELECT
            date_period,
            gc_platform,
            bgr,
            gc_city,
            brand,
            SUM(case when is_own_brand=1 then reported_offtake_mrp else normalized_offtake end) AS actual_sum,
            SUM(SUM(normalized_offtake)) 
              OVER (PARTITION BY date_period, gc_platform, bgr,gc_city)
              AS category_size,
            CASE 
              WHEN SUM(SUM(normalized_offtake)) 
                      OVER (PARTITION BY date_period, gc_platform, bgr,gc_city) = 0 
                THEN 0
              ELSE  
                ROUND(
                  100 * sum(normalized_offtake)
                        / SUM(SUM(normalized_offtake)) 
                            OVER (PARTITION BY date_period, gc_platform, bgr,gc_city), 
                  1
                )
            END AS market_share
          FROM filtered
          WHERE date_period IS NOT NULL
          GROUP BY date_period, gc_platform, bgr,gc_city,brand
          ORDER BY bgr,brand,gc_platform,date_period)
          ,is_own_brand as (select distinct gc_platform,brand,is_own_brand from filtered)
          select agg.*,is_own_brand.is_own_brand from agg left join is_own_brand on agg.gc_platform = is_own_brand.gc_platform  and agg.brand = is_own_brand.brand order by bgr,brand,gc_platform,date_period
          """
          result = sfFetch(query2)
          result.columns = result.columns.str.lower()
          result['account'] = db_name

      return result

  except Exception as e:
      print(f"[ERROR] account={db_name}: {e}", file=sys.stderr)
      return pd.DataFrame()




def compare_monthly_metrics(df, date_col='date_period', group_cols=None, metric_cols=None):
    
    if group_cols is None:
        group_cols = ['account','gc_platform', 'bgr','brand']
    
    if metric_cols is None:
        metric_cols = ['actual_sum','actual_sale','market_share']
    # else:
    #     metric_cols=[
    #         'wt_osa', 'wt_osa_in_ls', 'wt_discounting', 
    #         'sov', 'ad_sov', 'organic_sov'
    #     ]
    # Filter for min and max month
    df_min = df.loc[df[date_col] == 'dr2']
    df_max = df.loc[df[date_col] == 'dr1']

    # Merge on grouping columns, drop month columns
    df_merged = df_max.drop(columns=[date_col]).merge(
        df_min.drop(columns=[date_col]),
        on=group_cols,
        how='left',
        suffixes=('_dr1', '_dr2')
    )

    # Calculate changes
    for col in metric_cols:
        # Convert columns to numeric, coercing errors to NaN
        df_merged[f'{col}_dr1'] = pd.to_numeric(df_merged[f'{col}_dr1'], errors='coerce')
        df_merged[f'{col}_dr2'] = pd.to_numeric(df_merged[f'{col}_dr2'], errors='coerce')
        
        # Now perform the calculation
        df_merged[f'change_{col}'] = round(df_merged[f'{col}_dr1'] - df_merged[f'{col}_dr2'],1)

    # Select relevant columns to return
    selected_cols = group_cols + [f'{col}_dr1' for col in metric_cols] + [f'change_{col}' for col in metric_cols]
    return df_merged[selected_cols]




def get_ppu_cat_brand(
    acc,
    
    l1_start, l1_end,
    l2_start, l2_end,
    n_bins=3,
    clip_l=0.05,
    clip_u=0.95,
    random_state=42
):
    db_name = acc['db_name']
    query = f"""
    SELECT
    lower(qpm.gc_platform) as gc_platform,
      lower(qpd.bgr) AS bgr,
      qpm.brand,
      ROUND(
        SUM(CASE WHEN qpm.ppu_flag = '1'
                  THEN qpm.price_per_100_units * GREATEST(qpm.estimated_offtake_mrp,1)
                  ELSE 0 END)
        / NULLIF(SUM(CASE WHEN qpm.ppu_flag = '1'
                         THEN GREATEST(qpm.estimated_offtake_mrp,1)
                         ELSE 0 END),0),
        0
      ) AS wt_ppu
    FROM {db_name}.aggregate.qcom_pid_metrics qpm
    LEFT JOIN {db_name}.aggregate.qcom_product_dimensions qpd
      ON qpm.product_id = qpd.product_id
    WHERE  (qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}') 
    GROUP BY 1,2,3
    ORDER BY 1,2,3
    """
    br = sfFetch(query)
    # print(query)
    query2=f"""
    select distinct lower(gc_platform) as gc_platform,lower(bgr) as bgr,brand ,is_own_brand from {db_name}.aggregate.qcom_product_dimensions

    """
    br2=sfFetch(query2)
    br2.columns=br2.columns.str.lower()
    br.columns = br.columns.str.lower()

    br=br.merge(br2,on=['gc_platform','bgr','brand'],how='left')
    # return br
    br['account'] = db_name
    br=br.dropna(subset=['wt_ppu'])
    br['wt_ppu']=abs(br['wt_ppu'])
    kb_ppu  = KBinsDiscretizer(n_bins=n_bins, encode='ordinal', strategy='kmeans', random_state=random_state)
    # kb_gram = KBinsDiscretizer(n_bins=n_bins, encode='ordinal', strategy='kmeans', random_state=random_state)
    out = []

    for bgr, grp in br.groupby(['gc_platform','bgr'], as_index=False):
        sub = grp.copy()
        # PPU pipeline
        med_ppu = sub['wt_ppu'].median()
        sub['ppu_imp'] = sub['wt_ppu'].fillna(med_ppu)
        lo_ppu, hi_ppu = sub['ppu_imp'].quantile([clip_l, clip_u])
        sub['ppu_win'] = sub['ppu_imp'].clip(lo_ppu, hi_ppu)
        sub['log_ppu'] = np.log1p(sub['ppu_win'])

        # Grammage pipeline
        # med_g = sub['grammage'].median()
        # sub['gram_imp'] = sub['grammage'].fillna(med_g)
        # lo_g, hi_g = sub['gram_imp'].quantile([clip_l, clip_u])
        # sub['gram_win'] = sub['gram_imp'].clip(lo_g, hi_g)
        # sub['log_gram'] = np.log1p(sub['gram_win'])

        if len(sub) >= n_bins:
            sub['binning_method'] = 'kmeans'
      
            bins_ppu  = kb_ppu .fit_transform(sub[['log_ppu']]).astype(int).flatten()
            
            # bins_gram = kb_gram.fit_transform(sub[['log_gram']]).astype(int).flatten()
            sub['ppu_category']      = pd.Categorical.from_codes(bins_ppu, categories=['cost_effective','mid_range','premium'][:n_bins])
            # sub['grammage_category'] = pd.Categorical.from_codes(bins_gram, categories=['small','medium','large'][:n_bins])
        else:
            sub['binning_method'] = 'quantile'
            q1_ppu, q2_ppu = sub['wt_ppu'].quantile([1/3,2/3])
            sub['ppu_category'] = sub['wt_ppu'].apply(
                lambda x: 'cost_effective' if x <= q1_ppu else ('mid_range' if x <= q2_ppu else 'premium')
            )
            # q1_g, q2_g = sub['grammage'].quantile([1/3,2/3])
            # sub['grammage_category'] = sub['grammage'].apply(
            #     lambda x: 'small' if x <= q1_g else ('medium' if x <= q2_g else 'large')
            # )

        out.append(sub)
    df = pd.concat(out, ignore_index=True)

    ppu_ranges = (
        df
        .groupby(['gc_platform','bgr','ppu_category'], as_index=False)['wt_ppu']
        .agg(lower=lambda s: math.floor(s.min()/5)*5,
             upper=lambda s: math.ceil( s.max()/5)*5)
        .rename(columns={'lower':'ppu_lower_bound','upper':'ppu_upper_bound'})
    )
    # gram_ranges = (
    #     df
    #     .groupby(['gc_platform','bgr','grammage_category'], as_index=False)['grammage']
    #     .agg(lower=lambda s: math.floor(s.min()/5)*5,
    #          upper=lambda s: math.ceil( s.max()/5)*5)
    #     .rename(columns={'lower':'grammage_lower_bound','upper':'grammage_upper_bound'})
    # )

    df = df.merge(ppu_ranges, on=['gc_platform','bgr','ppu_category'], how='left')
    # df = df.merge(gram_ranges, on=['gc_platform','bgr','grammage_category'], how='left')

    return df.drop(columns=['ppu_imp','ppu_win','log_ppu'])
    # return df.drop(columns=['ppu_imp','ppu_win','log_ppu','gram_imp','gram_win','log_gram'])



def get_grammage(acc,start_date1,end_date1,start_date2,end_date2):
    db_name = acc['db_name']
    query = f"""
    SELECT
    lower(qpm.gc_platform) as gc_platform,
      lower(qpd.bgr) AS bgr,
      qpm.brand,
      qpm.product_id,
      AVG(qpd.grammage) AS grammage
    FROM {db_name}.aggregate.qcom_pid_metrics qpm
    LEFT JOIN {db_name}.aggregate.qcom_product_dimensions qpd
      ON qpm.product_id = qpd.product_id
    WHERE (
        qpm.snapshot_date BETWEEN DATE '{start_date1.date()}'
                            AND LEAST(
                                 DATEADD(month,-1,qpm.platform_offtake_uploaded_till),
                                 DATE '{end_date1.date()}'
                               )
      )
      OR (
        qpm.snapshot_date BETWEEN DATE '{start_date2.date()}'
                            AND LEAST(
                                 qpm.platform_offtake_uploaded_till,
                                 DATE '{end_date2.date()}'
                               )
      )
    GROUP BY 1,2,3,4
    ORDER BY 1,2,3,4
    """
    br = sfFetch(query)
    br.columns = br.columns.str.lower()
    br['account'] = db_name
    return br



def get_metrics(acc,l1_start,l1_end,l2_start,l2_end):
    db_name=acc['db_name']
    df1=sfFetch(f"""with valid_pids as (
    select distinct product_id
    from {db_name}.AGGREGATE.qcom_product_dimensions
    where is_platform_bundle is null
),

weighted_metrics as (
    select 
        qpm.gc_platform,
                CASE
                  WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}'
                    THEN 'dr1'
                  WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                        AND DATE '{l2_end.date()}'
                    THEN 'dr2'
                  ELSE NULL
                END AS date_period,
        lower(qpm.bgr) as bgr,
        qpm.gc_city,
        qpm.brand,
        qpm.total_available,
        qpm.total_listed,
        qpm.total_scrapes,
        (qpm.total_available * GREATEST(COALESCE(qpm.estimated_offtake_mrp_last_month, qpm.estimated_offtake_mrp_current_month), 1)) as wt_total_available,
        (qpm.total_listed * GREATEST(COALESCE(qpm.estimated_offtake_mrp_last_month, qpm.estimated_offtake_mrp_current_month), 1)) as wt_total_listed,
        (qpm.total_scrapes * GREATEST(COALESCE(qpm.estimated_offtake_mrp_last_month, qpm.estimated_offtake_mrp_current_month), 1)) as wt_total_scrapes
        -- (total_available * GREATEST(sku_importance, 1)) as wt_total_available,
        -- (total_listed * GREATEST(sku_importance, 1)) as wt_total_listed,
        -- (total_scrapes * GREATEST(sku_importance, 1)) as wt_total_scrapes
    from {db_name}.AGGREGATE.qcom_pid_metrics qpm
    join valid_pids vp
    on qpm.product_id = vp.product_id
    WHERE (
      qpm.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND 
                DATE '{l2_end.date()}'
              
    )
    OR
    (
      qpm.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND 
                DATE '{l1_end.date()}'
              
            )
      -- and qpm.pid_bgr_selection_rank = 1
),

wt_aggregated as (
    select 
        gc_platform,
        date_period,
        lower(bgr) as bgr,
        gc_city,
        brand,
        sum(wt_total_available) as total_wt_available,
        sum(wt_total_listed) as total_wt_listed,
        sum(wt_total_scrapes) as total_wt_scrapes
    from weighted_metrics
    group by 1,2,3,4,5
),

bgr_prices as (
    select 
        gc_platform,
         CASE
                  WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}'
                    THEN 'dr1'
                  WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                        AND DATE '{l2_end.date()}'
                    THEN 'dr2'
                  ELSE NULL
                END AS date_period,
        lower(bgr) as bgr,
        gc_city,
        brand,
        -- sum(mrp) as mrp,
        -- sum(price) as price,
        -- sum(mrp * GREATEST(sku_importance, 1)) as mrp,
        -- sum(price * GREATEST(sku_importance, 1)) as price,
        -- greatest(sum(estimated_qty_sold), 1) as qty
        sum(mrp * GREATEST(estimated_qty_sold, 1)) as wt_mrp,
        sum(price * GREATEST(estimated_qty_sold, 1)) as wt_price
    from {db_name}.AGGREGATE.qcom_pid_metrics qpm
    WHERE (
      qpm.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND  
                DATE '{l2_end.date()}'
              
    )
    OR
    (
      qpm.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND  
                DATE '{l1_end.date()}'
              
            )
      and asp_flag = 1
      -- and pid_bgr_selection_rank = 1
    group by 1,2,3,4,5
),

month_bgr_imps as (
    select 
        gc_platform, 
         CASE
                  WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}'
                    THEN 'dr1'
                  WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                        AND DATE '{l2_end.date()}'
                    THEN 'dr2'
                  ELSE NULL
                END AS date_period, 
        lower(bgr) as bgr,
        gc_city,
        brand,
        round(sum(total_impressions), 0) as ti, 
        round(sum(ad_impressions), 0) as ai, 
        round(sum(total_impressions) - sum(ad_impressions), 0) as oi
    from {db_name}.AGGREGATE.qcom_pid_metrics qpm
    WHERE (
      qpm.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND 
                DATE '{l2_end.date()}'
              
    )
    OR
    (
      qpm.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND 
                DATE '{l1_end.date()}'
              
            )
      -- and pid_bgr_selection_rank = 1
    group by 1,2,3,4,5
),

month_imps as (
    select 
        gc_platform, 
         CASE
                  WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}'
                    THEN 'dr1'
                  WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                        AND DATE '{l2_end.date()}'
                    THEN 'dr2'
                  ELSE NULL
                END AS date_period, 
        lower(bgr) as bgr, 
        gc_city,
        round(sum(total_impressions), 0) as ti,
        round(sum(ad_impressions), 0) as ai, 
        round(sum(total_impressions) - sum(ad_impressions), 0) as oi
    from {db_name}.AGGREGATE.qcom_pid_metrics qpm 
    WHERE (
      qpm.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND 
                DATE '{l2_end.date()}'
              
    )
    OR
    (
      qpm.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND 
                DATE '{l1_end.date()}'
              
            )
    -- and pid_bgr_selection_rank = 1
    group by 1,2,3,4
),

bgr_sov as (
    select 
        mbi.gc_platform,
        mbi.date_period,
        lower(mbi.bgr) as bgr,
        mbi.gc_city,
        mbi.brand,
        coalesce(round(mbi.ti::numeric / nullif(mi.ti, 0), 3) * 100.0, 0) as sov,
        coalesce(round(mbi.ai::numeric / nullif(mi.ai, 0), 3) * 100.0, 0) as ad_sov,
        coalesce(round(mbi.oi::numeric / nullif(mi.oi, 0), 3) * 100.0, 0) as organic_sov
    from month_bgr_imps mbi
    join month_imps mi 
      on mbi.gc_platform = mi.gc_platform 
     and mbi.date_period = mi.date_period 
     and mbi.bgr = mi.bgr and mbi.gc_city=mi.gc_city
),

merged as (
    select 
        wa.gc_platform,
        wa.date_period,
        wa.bgr,
        wa.gc_city,
        wa.brand,
        wa.total_wt_available,
        wa.total_wt_listed,
        wa.total_wt_scrapes,
        -- bp.mrp mrp,
        -- bp.price price,
        -- bp.qty qty,
        bp.wt_mrp,
        bp.wt_price,
        bs.sov,
        bs.ad_sov,
        bs.organic_sov
    from wt_aggregated wa
    left join bgr_prices bp
      on wa.gc_platform = bp.gc_platform 
     and wa.date_period = bp.date_period 
     and wa.bgr = bp.bgr and wa.gc_city=bp.gc_city and wa.brand=bp.brand
    left join bgr_sov bs
      on wa.gc_platform = bs.gc_platform 
     and wa.date_period = bs.date_period 
     and wa.bgr = bs.bgr and wa.gc_city=bs.gc_city and wa.brand=bs.brand
),
brand_info as (select distinct gc_platform,brand,is_own_brand from {db_name}.aggregate.qcom_pid_metrics)

select 
    lower(merged.gc_platform) gc_platform,
    merged.date_period,
    lower(merged.bgr) bgr,
    merged.gc_city,
    merged.brand,
    bi.is_own_brand,
    round(100.0 * total_wt_available / nullif(total_wt_scrapes, 0), 1) as wt_osa,
    round(100.0 * total_wt_available / nullif(total_wt_listed, 0), 1) as wt_osa_in_ls,
    round(100.0 * (wt_mrp - wt_price) / nullif(wt_mrp, 0), 1) as wt_discounting,
    sov,
    ad_sov,
    organic_sov
from merged
left join brand_info bi on merged.brand=bi.brand and merged.gc_platform=bi.gc_platform
order by 1,2,3,4,5;""")
    
    df1.columns=[col.lower() for col in df1.columns]

    # df1['date_period']=pd.to_datetime(df1['date_period']).dt.date_period
    l2=df1.loc[df1['date_period']=='dr2'] 
    l1=df1.loc[df1['date_period']=='dr1']

    l1=l1.drop(columns='date_period').merge(l2.drop(columns='date_period'),on=['gc_platform','bgr', 'gc_city', 'brand', 'is_own_brand'],how='left',suffixes=('_dr1','_dr2'))
    l1['change_sov']=l1['sov_dr1']-l1['sov_dr2']
    l1['change_ogsov']=l1['organic_sov_dr1']-l1['organic_sov_dr2']
    l1['change_adsov']=l1['ad_sov_dr1']-l1['ad_sov_dr2']
    l1['change_wt_osa']=l1['wt_osa_dr1']-l1['wt_osa_dr2']
    l1['change_wt_osa_listing']=l1['wt_osa_in_ls_dr1']-l1['wt_osa_in_ls_dr2']
    l1['change_wt_discounting']=l1['wt_discounting_dr1']-l1['wt_discounting_dr2']
    l1=l1[['gc_platform', 'bgr', 'gc_city', 'brand', 'is_own_brand', 'wt_osa_dr1',
        'wt_osa_in_ls_dr1', 'wt_discounting_dr1', 'sov_dr1', 'ad_sov_dr1',
        'organic_sov_dr1','change_sov', 'change_ogsov',
        'change_adsov', 'change_wt_osa', 'change_wt_osa_listing',
        'change_wt_discounting']]


    return l1
    


def get_bgr_metrics(acc,l1_start,l1_end,l2_start,l2_end):
    db_name=acc['db_name']
    query=f"""with valid_pids as (
    select distinct product_id
    from {db_name}.AGGREGATE.qcom_product_dimensions
    where is_platform_bundle is null
),

weighted_metrics as (
    select 
        qpm.gc_platform,
                CASE
                  WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}'
                    THEN 'dr1'
                  WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                        AND DATE '{l2_end.date()}'
                    THEN 'dr2'
                  ELSE NULL
                END AS date_period,
        lower(qpm.bgr) as bgr,
        qpm.brand,
        qpm.total_available,
        qpm.total_listed,
        qpm.total_scrapes,
        (qpm.total_available * GREATEST(COALESCE(qpm.estimated_offtake_mrp_last_month, qpm.estimated_offtake_mrp_current_month), 1)) as wt_total_available,
        (qpm.total_listed * GREATEST(COALESCE(qpm.estimated_offtake_mrp_last_month, qpm.estimated_offtake_mrp_current_month), 1)) as wt_total_listed,
        (qpm.total_scrapes * GREATEST(COALESCE(qpm.estimated_offtake_mrp_last_month, qpm.estimated_offtake_mrp_current_month), 1)) as wt_total_scrapes
        -- (total_available * GREATEST(sku_importance, 1)) as wt_total_available,
        -- (total_listed * GREATEST(sku_importance, 1)) as wt_total_listed,
        -- (total_scrapes * GREATEST(sku_importance, 1)) as wt_total_scrapes
    from {db_name}.AGGREGATE.qcom_pid_metrics qpm
    join valid_pids vp
    on qpm.product_id = vp.product_id
    WHERE (
      qpm.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND 
                DATE '{l2_end.date()}'
              
    )
    OR
    (
      qpm.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND 
                DATE '{l1_end.date()}'
              
            )
      -- and qpm.pid_bgr_selection_rank = 1
),

wt_aggregated as (
    select 
        gc_platform,
        date_period,
        lower(bgr) as bgr,
        brand,
        sum(wt_total_available) as total_wt_available,
        sum(wt_total_listed) as total_wt_listed,
        sum(wt_total_scrapes) as total_wt_scrapes
    from weighted_metrics
    group by 1,2,3,4
),

bgr_prices as (
    select 
        gc_platform,
         CASE
                  WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}'
                    THEN 'dr1'
                  WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                        AND DATE '{l2_end.date()}'
                    THEN 'dr2'
                  ELSE NULL
                END AS date_period,
        lower(bgr) as bgr,
        brand,
        -- sum(mrp) as mrp,
        -- sum(price) as price,
        -- sum(mrp * GREATEST(sku_importance, 1)) as mrp,
        -- sum(price * GREATEST(sku_importance, 1)) as price,
        -- greatest(sum(estimated_qty_sold), 1) as qty
        sum(mrp * GREATEST(estimated_qty_sold, 1)) as wt_mrp,
        sum(price * GREATEST(estimated_qty_sold, 1)) as wt_price
    from {db_name}.AGGREGATE.qcom_pid_metrics qpm
    WHERE ((
      qpm.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND  
                DATE '{l2_end.date()}'
              
    )
    OR
    (
      qpm.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND  
                DATE '{l1_end.date()}'
              
            ))
      and asp_flag = 1
      -- and pid_bgr_selection_rank = 1
    group by 1,2,3,4
),

month_bgr_imps as (
    select 
        gc_platform, 
         CASE
                  WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}'
                    THEN 'dr1'
                  WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                        AND DATE '{l2_end.date()}'
                    THEN 'dr2'
                  ELSE NULL
                END AS date_period, 
        lower(bgr) as bgr,
        brand,
        round(sum(total_impressions), 0) as ti, 
        round(sum(ad_impressions), 0) as ai, 
        round(sum(total_impressions) - sum(ad_impressions), 0) as oi
    from {db_name}.AGGREGATE.qcom_pid_metrics qpm
    WHERE (
      qpm.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND 
                DATE '{l2_end.date()}'
              
    )
    OR
    (
      qpm.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND 
                DATE '{l1_end.date()}'
              
            )
      -- and pid_bgr_selection_rank = 1
    group by 1,2,3,4
),

month_imps as (
    select 
        gc_platform, 
         CASE
                  WHEN qpm.snapshot_date BETWEEN DATE '{l1_start.date()}' 
                                        AND DATE '{l1_end.date()}'
                    THEN 'dr1'
                  WHEN qpm.snapshot_date BETWEEN DATE '{l2_start.date()}' 
                                        AND DATE '{l2_end.date()}'
                    THEN 'dr2'
                  ELSE NULL
                END AS date_period, 
        lower(bgr) as bgr, 
        round(sum(total_impressions), 0) as ti,
        round(sum(ad_impressions), 0) as ai, 
        round(sum(total_impressions) - sum(ad_impressions), 0) as oi
    from {db_name}.AGGREGATE.qcom_pid_metrics qpm 
    WHERE (
      qpm.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND 
                DATE '{l2_end.date()}'
              
    )
    OR
    (
      qpm.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND 
                DATE '{l1_end.date()}'
              
            )
    -- and pid_bgr_selection_rank = 1
    group by 1,2,3
),

bgr_sov as (
    select 
        mbi.gc_platform,
        mbi.date_period,
        lower(mbi.bgr) as bgr,
        mbi.brand,
        coalesce(round(mbi.ti::numeric / nullif(mi.ti, 0), 3) * 100.0, 0) as sov,
        coalesce(round(mbi.ai::numeric / nullif(mi.ai, 0), 3) * 100.0, 0) as ad_sov,
        coalesce(round(mbi.oi::numeric / nullif(mi.oi, 0), 3) * 100.0, 0) as organic_sov
    from month_bgr_imps mbi
    join month_imps mi 
      on mbi.gc_platform = mi.gc_platform 
     and mbi.date_period = mi.date_period 
     and mbi.bgr = mi.bgr 
),

merged as (
    select 
        wa.gc_platform,
        wa.date_period,
        wa.bgr,
        wa.brand,
        wa.total_wt_available,
        wa.total_wt_listed,
        wa.total_wt_scrapes,
        -- bp.mrp mrp,
        -- bp.price price,
        -- bp.qty qty,
        bp.wt_mrp,
        bp.wt_price,
        bs.sov,
        bs.ad_sov,
        bs.organic_sov
    from wt_aggregated wa
    left join bgr_prices bp
      on wa.gc_platform = bp.gc_platform 
     and wa.date_period = bp.date_period 
     and wa.bgr = bp.bgr and  wa.brand=bp.brand
    left join bgr_sov bs
      on wa.gc_platform = bs.gc_platform 
     and wa.date_period = bs.date_period 
     and wa.bgr = bs.bgr  and wa.brand=bs.brand
),
brand_info as (select distinct gc_platform,brand,is_own_brand from {db_name}.aggregate.qcom_pid_metrics)

select 
    lower(merged.gc_platform) gc_platform,
    merged.date_period,
    lower(merged.bgr) bgr,
    merged.brand,
    bi.is_own_brand,
    round(100.0 * total_wt_available / nullif(total_wt_scrapes, 0), 1) as wt_osa,
    round(100.0 * total_wt_available / nullif(total_wt_listed, 0), 1) as wt_osa_in_ls,
    round(100.0 * (wt_mrp - wt_price) / nullif(wt_mrp, 0), 1) as wt_discounting,
    sov,
    ad_sov,
    organic_sov
from merged
left join brand_info bi on merged.brand=bi.brand and merged.gc_platform=bi.gc_platform
order by 1,2,3,4,5;"""

    df1=sfFetch(query)
    df1.columns=[col.lower() for col in df1.columns]

    # df1['date_period']=pd.to_datetime(df1['date_period']).dt.date_period
    l2=df1.loc[df1['date_period']=='dr2'] 
    l1=df1.loc[df1['date_period']=='dr1']

    l1=l1.drop(columns='date_period').merge(l2.drop(columns='date_period'),on=['gc_platform','bgr', 'brand', 'is_own_brand'],how='left',suffixes=('_dr1','_dr2'))
    l1['change_sov']=l1['sov_dr1']-l1['sov_dr2']
    l1['change_ogsov']=l1['organic_sov_dr1']-l1['organic_sov_dr2']
    l1['change_adsov']=l1['ad_sov_dr1']-l1['ad_sov_dr2']
    l1['change_wt_osa']=l1['wt_osa_dr1']-l1['wt_osa_dr2']
    l1['change_wt_osa_listing']=l1['wt_osa_in_ls_dr1']-l1['wt_osa_in_ls_dr2']
    l1['change_wt_discounting']=l1['wt_discounting_dr1']-l1['wt_discounting_dr2']
    l1=l1[['gc_platform', 'bgr', 'brand', 'is_own_brand', 'wt_osa_dr1',
        'wt_osa_in_ls_dr1', 'wt_discounting_dr1', 'sov_dr1', 'ad_sov_dr1',
        'organic_sov_dr1','change_sov', 'change_ogsov',
        'change_adsov', 'change_wt_osa', 'change_wt_osa_listing',
        'change_wt_discounting']]


    return l1
    

def calculate_impact(ms,own):
    dr2 = ms[ms['date_period'] == 'dr2'] \
            .rename(columns={'category_size':'size_dr2','market_share':'share_dr2','actual_sum':'actual_sales_dr2'}) \
            .drop(columns=['date_period','is_own_brand'])
    dr1 = ms[ms['date_period'] == 'dr1'] \
            .rename(columns={'category_size':'size_dr1','market_share':'share_dr1','actual_sum':'actual_sales_dr1'}) \
            .drop(columns=['date_period','is_own_brand'])
    if own==1:
        merge_cols = ['account','gc_platform','bgr']
    elif own==2:
        merge_cols = ['account','gc_platform','bgr','brand']
    elif own==3:
        merge_cols = ['account','gc_platform','bgr','gc_city']
    else:
        merge_cols = ['account','gc_platform','bgr','gc_city','brand']
    merged = dr1.merge(dr2, on=merge_cols, how='inner')

    merged['pct_change']    = merged['share_dr1'] - merged['share_dr2']
    merged['sales_impact']  = (merged['pct_change'] * merged['size_dr1'])/100    
    return merged


def get_product_level_data(acc, l2_start, l2_end, l1_start, l1_end, platform):
    db_name = acc['db_name']
    query = f"""
    WITH filtered AS (
      SELECT
        -- label each row as l1week or l2week
        
        metrics.product_id,
        mapping.item_id        AS item_id,
        metrics.title          AS title,
        metrics.original_grammage,
        lower(metrics.bgr)     AS bgr,
        metrics.gc_city,
        metrics.brand,
        CASE
          WHEN snapshot_date BETWEEN DATE '{l2_start.date()}' AND DATE '{l2_end.date()}' THEN 'dr2'
          WHEN snapshot_date BETWEEN DATE '{l1_start.date()}' AND DATE '{l1_end.date()}' THEN 'dr1'
          ELSE NULL
        END AS date_period,
        metrics.reported_offtake_mrp,
        metrics.reported_offtake_sp,
        metrics.price,
        metrics.mrp,
        metrics.reported_offtake_qty,
        metrics.total_available,
        metrics.total_listed,
        metrics.total_scrapes,
        metrics.estimated_offtake_mrp_last_month,
        metrics.estimated_offtake_mrp_current_month,
        metrics.estimated_qty_sold,
        metrics.asp_flag,
        metrics.normalized_offtake,
        metrics.total_impressions,
        metrics.organic_impressions,
        metrics.ad_impressions,
        metrics.ppu_flag,
        metrics.price_per_100_units
      FROM {db_name}.aggregate.{platform}_qcom_pid_metrics AS metrics
      LEFT JOIN {db_name}.aggregate.{platform}_item_pid_mapping AS mapping
        ON metrics.product_id = mapping.product_id
       AND mapping.type = 'single'
       AND mapping.gc_platform = metrics.gc_platform
      WHERE (
      metrics.snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND LEAST(
                dateadd(month,-1,metrics.platform_offtake_uploaded_till), 
                DATE '{l2_end.date()}'
              )
    )
    OR
    (
      metrics.snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND LEAST(
                metrics.platform_offtake_uploaded_till, 
                DATE '{l1_end.date()}'
              )
            )
        AND metrics.is_own_brand IN (0,1)
    )
    SELECT
      
      product_id           AS "product_id",
      item_id              AS "item_id",
      title,
      original_grammage    AS "grammage",
      bgr,
      gc_city,
      brand,
      date_period,
      SUM(reported_offtake_mrp)           AS "Offtake MRP",
      SUM(reported_offtake_sp)            AS "Offtake SP",
      AVG(CASE WHEN asp_flag='1' THEN price END) AS "SP",
      AVG(CASE WHEN asp_flag='1' THEN mrp   END) AS "MRP",
      SUM(reported_offtake_qty)           AS "Units",
      100*sum(total_listed)/sum(total_scrapes)                   AS "ds_listing",
      -- weighted availability
      SUM(total_available * GREATEST(COALESCE(estimated_offtake_mrp_last_month, estimated_offtake_mrp_current_month),1))
        * 100.0
        / NULLIF(SUM(total_scrapes * GREATEST(COALESCE(estimated_offtake_mrp_last_month, estimated_offtake_mrp_current_month),1)),0)
        AS "wt_osa",
        SUM(total_available * GREATEST(COALESCE(estimated_offtake_mrp_last_month, estimated_offtake_mrp_current_month),1))
        * 100.0
        / NULLIF(SUM(total_listed * GREATEST(COALESCE(estimated_offtake_mrp_last_month, estimated_offtake_mrp_current_month),1)),0)
        AS "wt_osa_in_ls",
      -- weighted discount
      (
        SUM(CASE WHEN asp_flag='1' THEN mrp * GREATEST(estimated_qty_sold,1) END)
        - SUM(CASE WHEN asp_flag='1' THEN price * GREATEST(estimated_qty_sold,1) END)
      ) * 100.0
        / NULLIF(SUM(CASE WHEN asp_flag='1' THEN mrp * GREATEST(estimated_qty_sold,1) END),0)
        AS "wt_discounting",
      -- category share within that date_period
      SUM(normalized_offtake)
        / NULLIF(
            SUM(SUM(normalized_offtake)) OVER (PARTITION BY bgr, gc_city, date_period),
            0
        ) * 100.0 AS "market_share",
      SUM(total_impressions)
        * 100.0
        / NULLIF(
            SUM(SUM(total_impressions)) OVER (PARTITION BY bgr, gc_city, date_period),
            0
        ) AS "sov",
      SUM(organic_impressions)
        * 100.0
        / NULLIF(
            SUM(SUM(organic_impressions)) OVER (PARTITION BY bgr, gc_city, date_period),
            0
        ) AS "organic_sov",
      SUM(ad_impressions)
        * 100.0
        / NULLIF(
            SUM(SUM(ad_impressions)) OVER (PARTITION BY bgr, gc_city, date_period),
            0
        ) AS "ad_sov",
      SUM(
        CASE WHEN ppu_flag='1'
             THEN price_per_100_units * GREATEST(estimated_offtake_mrp_last_month,estimated_offtake_mrp_current_month)
        END
      )
        / NULLIF(
            SUM(CASE WHEN ppu_flag='1'
                     THEN GREATEST(estimated_offtake_mrp_last_month,estimated_offtake_mrp_current_month)
                END),
            0
        ) AS "Wt. PPU (x100)"
    FROM filtered
    WHERE date_period IS NOT NULL
    GROUP BY
      1,2,3,4,5,6,7,8
    
    
    """
    df = sfFetch(query)
    return df

# Add these imports at the top of your notebook
from concurrent.futures import ThreadPoolExecutor, as_completed
import threading
# Wrapper function for each platform's data collection
def get_platform_data(acc, l2_start, l2_end, l1_start, l1_end, platform):
    try:
        # print(f"Starting data collection for {platform}")
        df = get_product_level_data(acc, l2_start, l2_end, l1_start, l1_end, platform)
        # print(f"Completed data collection for {platform}")
        return platform, df
    except Exception as e:
        # print(f"Error collecting data for {platform}: {str(e)}")
        return platform, None


def get_all_platforms_data_parallel(acc, l2_start, l2_end, l1_start, l1_end, platforms):
    all_dataframes = []

    with ThreadPoolExecutor(max_workers=len(platforms)) as executor:
        future_to_platform = {
            executor.submit(get_platform_data, acc, l2_start, l2_end, l1_start, l1_end, platform): platform
            for platform in platforms
        }

        for future in as_completed(future_to_platform):
            platform = future_to_platform[future]
            try:
                platform_name, df = future.result()
                if df is not None and not df.empty:
                    df['gc_platform'] = platform_name
                    all_dataframes.append(df)
                    # Successfully collected data
                    pass
                else:
                    # Failed to collect data
                    pass
            except Exception as e:
                # Exception occurred - will be handled by caller
                pass

    if all_dataframes:
        combined_df = pd.concat(all_dataframes, ignore_index=True)
        combined_df['account']=acc['db_name']
        combined_df.columns=combined_df.columns.str.lower()
		
        return combined_df
    else:
        return pd.DataFrame()  # Return empty DataFrame if nothing succeeded



# Example usage


def get_bgr_ms(acc,l1_start,l1_end,l2_start,l2_end):

  query=f''' select '{acc['db_name']}' as account,CASE
          WHEN snapshot_date BETWEEN DATE '{l2_start.date()}' AND DATE '{l2_end.date()}' THEN 'dr2'
          WHEN snapshot_date BETWEEN DATE '{l1_start.date()}' AND DATE '{l1_end.date()}' THEN 'dr1'
          ELSE NULL
        END AS date_period,gc_platform,bgr,brand,is_own_brand,
        SUM(case when is_own_brand=1 then reported_offtake_mrp else normalized_offtake end) AS actual_sum,
              SUM(SUM(normalized_offtake)) 
                OVER (PARTITION BY date_period, gc_platform, bgr)
                AS category_size,
              CASE 
                WHEN SUM(SUM(normalized_offtake)) 
                       OVER (PARTITION BY date_period, gc_platform, bgr) = 0 
                  THEN 0
                ELSE  
                  ROUND(
                    100 * sum(normalized_offtake)
                          / SUM(SUM(normalized_offtake)) 
                              OVER (PARTITION BY date_period, gc_platform, bgr), 
                    1
                  )
              END AS market_share,
              from {acc['db_name']}.aggregate.qcom_pid_metrics 
        where (snapshot_date 
      BETWEEN DATE '{l2_start.date()}' 
          AND LEAST(
                dateadd(month,-1,platform_offtake_uploaded_till), 
                DATE '{l2_end.date()}'
              )
    )
    OR
    (
      
  snapshot_date 
      BETWEEN DATE '{l1_start.date()}' 
          AND LEAST(
                platform_offtake_uploaded_till, 
                DATE '{l1_end.date()}'
              )
            )
        AND is_own_brand IN (0,1)
    
    group by 1,2,3,4 ,5,6'''
  bgr_ms=sfFetch(query)
  type(bgr_ms)

  bgr_ms.columns=bgr_ms.columns.str.lower()
  bgr_ms['gc_platform']=bgr_ms['gc_platform'].str.lower()
  bgr_ms['bgr']=bgr_ms['bgr'].str.lower()
  return bgr_ms


def get_full_data(acc):
    br=get_ppu_cat_brand(acc,l1_start, l1_end,l2_start,l2_end)
    ms=fetch_for_acc(acc,l1_start, l1_end,l2_start,l2_end)

    merged_city=calculate_impact(ms[ms['is_own_brand']==1],4)
    merged_city['is_own_brand']=1
    merged_brand=calculate_impact(ms[ms['is_own_brand']==1],4)
    merged_brand['is_own_brand']=1
    bgr_ms=get_bgr_ms(acc,l1_start,l1_end,l2_start,l2_end)
    merged_category=calculate_impact(bgr_ms[bgr_ms['is_own_brand']==1],2)
    merged_category['is_own_brand']=1
    
    merged_comp=calculate_impact(ms[ms['is_own_brand']==0],own=4)
    merged_comp['is_own_brand']=0
    # print(merged_comp.head(10).to_markdown())
    merged_brand=merged_brand.merge(br,on=['account','gc_platform','bgr','brand','is_own_brand'],how='left')
    merged_comp=merged_comp.merge(br,on=['account','gc_platform','bgr','brand'],how='left')

   
    merged_city=merged_city[merged_city['gc_city']!='Others']
    merged_brand=merged_brand[merged_brand['gc_city']!='Others']
    merged_comp=merged_comp[merged_comp['gc_city']!='Others']
    merged_category1=merged_city.query('sales_impact<0').groupby(['account','gc_platform','bgr','brand','is_own_brand']).agg({'sales_impact':'sum'}).reset_index()
    merged_category_above=merged_city.query('sales_impact>0').groupby(['account','gc_platform','bgr','brand','is_own_brand']).agg({'sales_impact':'sum'}).reset_index()
    merged_category_top=merged_category_above.merge(merged_category.drop(columns='sales_impact'),on=['account','gc_platform','bgr','brand','is_own_brand'],how='left')
    merged_category=merged_category1.merge(merged_category.drop(columns='sales_impact'),on=['account','gc_platform','bgr','brand','is_own_brand'],how='left')
    
    
    below2 = (merged_category[merged_category['sales_impact']<-100000].sort_values(by=['account','gc_platform','sales_impact'], ascending=[True,True,True])
        .groupby(['account','gc_platform'])
        .head(2)
        .reset_index(drop=True))
    top2=(merged_category_top[merged_category_top['sales_impact']>50000].sort_values(['account','gc_platform','sales_impact'], ascending=[True,True,False])
        .groupby(['account','gc_platform'])
        .head(2)
        .reset_index(drop=True)
    )
    # print(below2.to_markdown())
    l1=get_metrics(acc,l1_start,l1_end,l2_start,l2_end)
    merged_comp=merged_comp.merge(l1,on=['gc_platform','bgr','gc_city','brand'],how='left')
    platforms = ['blinkit', 'instamart', 'zepto']
    all_data = get_all_platforms_data_parallel(acc, l2_start, l2_end, l1_start, l1_end, platforms)
    a=compare_monthly_metrics(all_data,date_col='date_period',group_cols=['account','gc_platform','bgr','gc_city','brand','product_id','title'],metric_cols=['ds_listing','wt_osa','wt_osa_in_ls','wt_discounting','sov','ad_sov','organic_sov','market_share'])
    a=a.merge(br,on=['account','gc_platform','brand','bgr'],how='left')
    a=a.merge(get_grammage(acc,l2_start,l2_end,l1_start,l1_end),on=['account','gc_platform','brand','bgr','product_id'],how='left')
    filter_comp=pd.DataFrame()
    sku_own=pd.DataFrame()
    city_df= pd.DataFrame()
    full_city_df=pd.DataFrame()
    city_df_brand=pd.DataFrame()
    for index,row in below2.iterrows():
        
        df_bad=merged_city[merged_city['sales_impact']<-10000].loc[(merged_city['gc_platform']==row['gc_platform']) & (merged_city['bgr']==row['bgr']) & (merged_city['brand']==row['brand'])].sort_values(by='sales_impact',ascending=True)

        if df_bad.empty:
            continue
        
        city_df=pd.concat([city_df,df_bad.head(1)])
        bad_df=df_bad.head(3)
        df_bad=df_bad.head(1)
        full_city_df=pd.concat([full_city_df,bad_df])
        gc_city=df_bad['gc_city'].values[0]
        df_brand=merged_brand[merged_brand['sales_impact']<-10000].loc[(merged_brand['gc_platform']==row['gc_platform']) & (merged_brand['bgr']==row['bgr']) & (merged_brand['gc_city']==gc_city) & (merged_brand['brand']==row['brand'])].sort_values(by='sales_impact',ascending=True).head(1)
        if df_brand.empty:
            continue
        ppu=df_brand['ppu_category'].values[0]


        
        city_df_brand=pd.concat([city_df_brand,df_brand])
        brand=df_brand['brand'].values[0]

        d2=merged_comp.loc[(merged_comp['gc_platform']==row['gc_platform']) & (merged_comp['bgr']==row['bgr']) & (merged_comp['gc_city']==gc_city) & (merged_comp['ppu_category']==ppu)]
        d2=d2[d2['pct_change']>0].sort_values(by='pct_change',ascending=False).head(2)
        d4=a.loc[(a['gc_platform']==row['gc_platform']) & (a['bgr']==row['bgr']) & (a['gc_city']==gc_city) & (a['brand']==brand)]
        d4=d4.sort_values(by='change_market_share',ascending=True).head(1)
        sku_own=pd.concat([sku_own,d4])
        filter_comp=pd.concat([filter_comp,d2])
        # print("filter_comp",filter_comp.to_markdown())

    # if city_df_brand.empty:
    #     return None
    # if filter_comp.empty:
    #     for index,row in city_df_brand.iterrows():
    #         d2=merged_comp.loc[(merged_comp['gc_platform']==row['gc_platform']) & (merged_comp['bgr']==row['bgr']) & (merged_comp['gc_city']==row['gc_city'])]
    #         d2=d2[d2['pct_change']>0].sort_values(by='pct_change',ascending=False).head(2)
    #         filter_comp=pd.concat([filter_comp,d2])


    # --- Fallback if no comps found yet ---
    if filter_comp.empty:
        if not city_df_brand.empty:
            tmp_list = []
            for _, row in city_df_brand.iterrows():
                d2 = merged_comp.loc[
                    (merged_comp['gc_platform'] == row['gc_platform']) &
                    (merged_comp['bgr']         == row['bgr']) &
                    (merged_comp['gc_city']     == row['gc_city'])
                ]
                d2 = d2.loc[d2['pct_change'] > 0].sort_values('pct_change', ascending=False).head(2)
                if not d2.empty:
                    tmp_list.append(d2)
            if tmp_list:
                filter_comp = pd.concat(tmp_list, ignore_index=True)
        # If still empty, consider a broader fallback (optional):
        if filter_comp.empty:
            # Example: take top movers per platform/bgr regardless of city
            fallback = (
                merged_comp.loc[merged_comp['pct_change'] > 0]
                            .sort_values(['gc_platform','bgr','pct_change'], ascending=[True,True,False])
                            .groupby(['gc_platform','bgr'], as_index=False)
                            .head(2)
            )
            if not fallback.empty:
                filter_comp = fallback.reset_index(drop=True)

    sku_comp = pd.DataFrame()
    if not filter_comp.empty:
        tmp_list = []
        for _, row in filter_comp.iterrows():
            d3 = a.loc[
                (a['gc_platform'] == row['gc_platform']) &
                (a['bgr']         == row['bgr']) &
                (a['gc_city']     == row['gc_city']) &
                (a['brand']       == row['brand'])
            ].sort_values('change_market_share', ascending=False).head(1)
            if not d3.empty:
                tmp_list.append(d3)
        if tmp_list:
            sku_comp = pd.concat(tmp_list, ignore_index=True)

    if not city_df_brand.empty:
        keys = ['gc_platform','bgr','gc_city','brand','is_own_brand']
        missing_in_city = [k for k in keys if k not in city_df_brand.columns]
        missing_in_l1   = [k for k in keys if k not in l1.columns]
        if not missing_in_city and not missing_in_l1:
            city_df_brand = city_df_brand.merge(l1, on=keys, how='left')
        else:
            pass





    # print(filter_comp.to_markdown())
    

    # filter_comp=filter_comp.merge(l1,on=['gc_platform','bgr','gc_city','brand','is_own_brand'],how='left')

    
    # sku_comp=pd.DataFrame()
    # for index,row in filter_comp.iterrows():
    #     # ppu=sku_own.loc[((sku_own['gc_platform']==row['gc_platform']) & (sku_own['bgr']==row['bgr']) & (sku_own['gc_city']==row['gc_city'])),'ppu_category' ].values
    #     # print(ppu)
    #     d3=a.loc[(a['gc_platform']==row['gc_platform']) & (a['bgr']==row['bgr']) & (a['gc_city']==row['gc_city']) & (a['brand']==row['brand'])]
    #     # d3=d3[d3['ppu_category'].isin(ppu)]
    #     d3=d3.sort_values(by='change_market_share',ascending=False).head(1)
    #     sku_comp=pd.concat([sku_comp,d3])
   
    # city_df_brand=city_df_brand.merge(l1,on=['gc_platform','bgr','gc_city','brand','is_own_brand'],how='left')


    
    
        # --- Process top2 exactly like we did below2 ---
    top_city_df       = pd.DataFrame()
    full_city_df_top  = pd.DataFrame()
    city_df_brand_top = pd.DataFrame()
    filter_comp_top   = pd.DataFrame()
    sku_own_top       = pd.DataFrame()
    sku_comp_top      = pd.DataFrame()

    for _, row in top2.iterrows():
        df_good = (merged_city[(merged_city['gc_platform'] == row['gc_platform']) &(merged_city['bgr']== row['bgr']) &(merged_city['sales_impact'] > 10000) & (merged_city['brand']==row['brand'])].sort_values(by='sales_impact', ascending=False))
        if df_good.empty:
            continue

        # Top city
        top_city_df      = pd.concat([top_city_df,     df_good.head(1)])
        full_city_df_top = pd.concat([full_city_df_top, df_good.head(3)])

        gc_city = df_good.iloc[0]['gc_city']

        # Now drill into brands within that city
        df_brand = (
            merged_brand[(merged_brand['gc_platform'] == row['gc_platform']) &
                         (merged_brand['bgr']         == row['bgr']) &
                         (merged_brand['gc_city']     == gc_city) & (merged_brand['brand']==row['brand']) &
                         (merged_brand['sales_impact']> 10000)]
            .sort_values(by='sales_impact', ascending=False)
            .head(1)
        )
        if df_brand.empty:
            continue

        city_df_brand_top = pd.concat([city_df_brand_top, df_brand])
        ppu   = df_brand.iloc[0]['ppu_category']
        brand = df_brand.iloc[0]['brand']

        # Competitive set for that city/ppu
        d2 = (
            merged_comp[(merged_comp['gc_platform'] == row['gc_platform']) &
                        (merged_comp['bgr']         == row['bgr']) &
                        (merged_comp['gc_city']     == gc_city) & 
                        (merged_comp['ppu_category']== ppu)]
            .loc[lambda df: df['pct_change'] > 0]
            .sort_values(by='pct_change', ascending=False)
            .head(2)
        )
        filter_comp_top = pd.concat([filter_comp_top, d2])

        # Own SKU driving the gain
        d4 = (
            a[(a['gc_platform'] == row['gc_platform']) &
              (a['bgr']         == row['bgr']) &
              (a['gc_city']     == gc_city) &
              (a['brand']       == brand)]
            .sort_values(by='change_market_share', ascending=False)
            .head(1)
        )
        sku_own_top = pd.concat([sku_own_top, d4])

    # If no comps found, fall back to top city brands
    if city_df_brand_top.empty and not top_city_df.empty:
        for _, row in city_df_brand_top.iterrows():
            d2 = (
                merged_comp[(merged_comp['gc_platform']==row['gc_platform']) &
                            (merged_comp['bgr']        ==row['bgr']) &
                            (merged_comp['gc_city']    ==row['gc_city'])]
                .loc[lambda df: df['pct_change'] > 0]
                .sort_values(by='pct_change', ascending=False)
                .head(2)
            )
            filter_comp_top = pd.concat([filter_comp_top, d2])

    # Enrich with L1 metrics
    # filter_comp_top = filter_comp_top.merge(
    #     l1,
    #     on=['gc_platform','bgr','gc_city','brand','is_own_brand'],
    #     how='left'
    # )

    # Pick the competitive SKU most responsible for the shift
    for _, row in filter_comp_top.iterrows():
        d3 = (
            a[(a['gc_platform'] == row['gc_platform']) &
              (a['bgr']         == row['bgr']) &
              (a['gc_city']     == row['gc_city']) &
              (a['brand']       == row['brand'])]
            .sort_values(by='change_market_share', ascending=False)
            .head(1)
        )
        sku_comp_top = pd.concat([sku_comp_top, d3])
    # Final merge to get the “full_top” DataFrame
    if not city_df_brand_top.empty:
        city_df_brand_top = city_df_brand_top.merge(
            l1,
            on=['gc_platform','bgr','gc_city','brand','is_own_brand'],
            how='left'
        )
    else:
        city_df_brand_top=pd.DataFrame()
    # full_top = (
    #     city_df_brand_top
    #     .merge(filter_comp_top, on=['gc_platform','bgr','gc_city','is_own_brand'], how='left', suffixes=('_brand','_comp'))
    #     .rename(columns={'account_brand':'account','brand_brand':'own_brand'})
    #     .drop(columns='account_comp')
    # )

    # Add your new objects into the return dict
    
    
    dict1={'ppu_cat':br,'metric':l1,'top2':top2,'below2':below2,'ms':ms,'filter_comp':filter_comp,'sku_own':sku_own,'sku_comp':sku_comp,'a':a,'city_df':city_df,'full_city_df':full_city_df,'city_df_brand':city_df_brand,'merged_category':merged_category,'merged_city':merged_city,'merged_brand':merged_brand,'merged_comp':merged_comp,'merged_category_top':merged_category_top,'merged_category_above':merged_category_above}
    dict1.update({
        'top_city_df':       top_city_df,
        'full_city_df_top':  full_city_df_top,
        'city_df_brand_top': city_df_brand_top,
        'filter_comp_top':   filter_comp_top,
        'sku_own_top':       sku_own_top,
        'sku_comp_top':      sku_comp_top,
        # 'full_top':          full_top,
    })
    return dict1


def get_bgr_data(acc,l1_start,l1_end,l2_start,l2_end):
    bgr_data=get_bgr_metrics(acc,l1_start,l1_end,l2_start,l2_end)
    bgr_data.columns=bgr_data.columns.str.lower()
    bgr_ms=get_bgr_ms(acc,l1_start,l1_end,l2_start,l2_end)

    bgr_ms=compare_monthly_metrics(bgr_ms,group_cols=['gc_platform','bgr','brand'],metric_cols=['market_share'])
    bgr_data=bgr_data.merge(bgr_ms,on=['gc_platform','bgr','brand'],how='left')
    return bgr_data

def get_sku_osa_data(acc, l1_start, l1_end, l2_start, l2_end, platform_update_date):

  db = acc['db_name']
  dfs = []

  for _, row in platform_update_date.iterrows():

    platform = row['GC_PLATFORM']
      # get plain YYYY-MM-DD strings
    lu = row['LAST_UPDATED'].isoformat()
    l1s = l1_start.date().isoformat()
    l1e = l1_end.date().isoformat()
    l2e = l2_end.date().isoformat()
    l2s = l2_start.date().isoformat()

    query = f"""
WITH osa_windows AS (
  SELECT
    gc_platform,
    gc_city,
    bgr,
    product_id,
    title,
    store_id,
    locality,

    CASE
      WHEN snapshot_date BETWEEN
           
             DATE '{l2s}'
           
        AND  DATE '{l2e}'
      THEN 'dr2'

      WHEN snapshot_date BETWEEN
              DATE '{l1s}'
        AND DATE '{l1e}'
      THEN 'dr1'
    END AS date_period,

    SUM(
      total_available
      * GREATEST(
          COALESCE(
            estimated_offtake_mrp_last_month,
            estimated_offtake_mrp_current_month
          ),
          1
        )
    ) * 100.0
      / NULLIF(
          SUM(
            total_scrapes
            * GREATEST(
                COALESCE(
                  estimated_offtake_mrp_last_month,
                  estimated_offtake_mrp_current_month
                ),
                1
              )
          ),
          0
        ) AS wt_osa

  FROM {db}.aggregate.qcom_locality_metrics
  WHERE
    (
      snapshot_date BETWEEN
        
           DATE '{l2s}'
        
      AND  DATE '{l2e}'
    )
    OR
    (
      snapshot_date BETWEEN
       
          DATE '{l1s}'
        
      AND  DATE '{l1e}'
    )
    AND gc_platform = '{platform}'
  GROUP BY 1,2,3,4,5,6,7,8
),
osa_data as(SELECT
  m.gc_platform,
  m.gc_city,
  m.bgr,
  m.product_id,
  m.title,
  m.store_id,
  m.locality,

  -- pull out the two windows side-by-side
  ROUND(l.wt_osa,  1) AS wt_osa_dr2,
  ROUND(m.wt_osa,  1) AS wt_osa_dr1,

  -- absolute and relative change
  ROUND(m.wt_osa - l.wt_osa,        1) AS osa_diff

FROM
  -- self-join LMTD → m  
  osa_windows l
  JOIN osa_windows m
    ON l.gc_platform = m.gc_platform
   AND l.gc_city     = m.gc_city
   AND l.bgr         = m.bgr
   AND l.product_id  = m.product_id
   AND l.title       = m.title
   AND l.locality    = m.locality
   AND l.store_id    = m.store_id

   AND l.date_period = 'dr2'
   AND m.date_period = 'dr1'
ORDER BY
  m.gc_platform,
  m.gc_city,
  m.bgr,
  m.product_id,
  m.locality
)
,
mtd_impact_pf AS (
  SELECT
  CASE
      WHEN snapshot_date BETWEEN
           
             DATE '{l2s}'
           
        AND  DATE '{l2e}'
      THEN 'dr2'

      WHEN snapshot_date BETWEEN
              DATE '{l1s}'
        AND DATE '{l1e}'
      THEN 'dr1'
    END AS date_period,
    gc_platform,
    gc_city,
    bgr,
    product_id,
    title,
    store_id,
    locality,
    SUM(total_available)       AS tot_ls,
    SUM(impact_scrore_for_psl_tags_pf) AS total_mtd_impact_pf
  FROM {db}.aggregate.qcom_locality_metrics
  WHERE
    (
      snapshot_date BETWEEN
       
           DATE '{l2s}'
        
      AND  DATE '{l2e}'
    )
    OR
    (
      snapshot_date BETWEEN
        
          DATE '{l1s}'
        
      AND  DATE '{l1e}'
    )
    AND gc_platform = '{platform}'
  GROUP BY 1,2,3,4,5,6,7,8
),
ranked_pf AS (
  SELECT
    *,
    SUM(total_mtd_impact_pf) OVER (PARTITION BY date_period,gc_platform, gc_city, bgr)
      AS grp_mtd_impact_pf,
    ROW_NUMBER() OVER (
      PARTITION BY date_period,gc_platform, gc_city, bgr
      ORDER BY total_mtd_impact_pf DESC
    ) AS sku_rank_pf,
    SUM(total_mtd_impact_pf) OVER (
      PARTITION BY date_period,gc_platform, gc_city, bgr
      ORDER BY total_mtd_impact_pf DESC
      ROWS BETWEEN UNBOUNDED PRECEDING AND CURRENT ROW
    )
    / NULLIF(
        SUM(total_mtd_impact_pf)
          OVER (PARTITION BY date_period,gc_platform, gc_city, bgr),
        0
      ) AS contribution_weight_pf
  FROM mtd_impact_pf
),

t AS (
  SELECT
    date_period,
    gc_platform,
    gc_city,
    lower(bgr) as bgr,
    product_id,
    title,
    store_id,
    locality,
    total_mtd_impact_pf,
    ROUND(contribution_weight_pf, 3) AS contribution_weight_pf,
    sku_rank_pf,
    CASE
      WHEN (contribution_weight_pf <= 0.25 OR sku_rank_pf <= 2)
           AND total_mtd_impact_pf > 0 THEN 'high'
      WHEN (contribution_weight_pf <= 0.50 OR sku_rank_pf <= 4)
           AND total_mtd_impact_pf > 0 THEN 'medium'
      ELSE 'low'
    END AS psl_tags_pf
  FROM ranked_pf
  ORDER BY bgr, gc_platform, gc_city, product_id, sku_rank_pf
)
SELECT
  t.*,
  h.total_listed,
  o.wt_osa_dr2,
  o.wt_osa_dr1,
  o.osa_diff
FROM t
LEFT JOIN (
  SELECT
  CASE
      WHEN snapshot_date =
           
        DATE '{l2e}'
      THEN 'dr2'

      WHEN snapshot_date= DATE '{l1e}'
      THEN 'dr1'
    END AS date_period,
    gc_platform,
    product_id,
    store_id,
    locality,
    total_listed
  FROM {db}.aggregate.qcom_locality_metrics
  WHERE
    snapshot_date = DATE '{l1e}' or snapshot_date = DATE '{l2e}'
    AND gc_platform = '{platform}'
) h
on t.date_period = h.date_period and
  t.gc_platform = h.gc_platform 
  AND t.product_id  = h.product_id
  AND t.locality    = h.locality
  AND t.store_id    = h.store_id
LEFT JOIN osa_data o

  ON 
  o.gc_platform = t.gc_platform
  AND o.product_id = t.product_id
  AND o.locality   = t.locality
  AND o.store_id    = t.store_id
WHERE
osa_diff < 0
;
""" 
    df = sfFetch(query)
    df.columns = df.columns.str.lower()
    dfs.append(df)

  return pd.concat(dfs, ignore_index=True)


def get_osa_summary(d,f,acc,l1_start,l1_end,l2_start,l2_end):
    bgr_data=get_bgr_data(acc,l1_start,l1_end,l2_start,l2_end)
    
    own_bgr_data=bgr_data[bgr_data['is_own_brand']==1]
    below2=d['below2'].sort_values(by='sales_impact',ascending=True).merge(own_bgr_data,on=['gc_platform','bgr','brand','is_own_brand'],how='left')
    osa_tab=d['city_df_brand'].sort_values(by='sales_impact',ascending=True).query('change_wt_osa<0')
    if osa_tab.empty:
        print('OSA not dropped')
        return pd.DataFrame(columns=['gc_platform','bgr','top_geo','share_dr1','pct_change','wt_osa_dr1','change_wt_osa','what_happened','potential_sales_loss','recommendation'])
    bgr_psl=f.groupby(['gc_platform','bgr'])['psl_mrp'].sum().round(0).reset_index()
    for index, row in osa_tab.iterrows():
        matching_row = bgr_psl.loc[
            (bgr_psl['gc_platform'] == row['gc_platform']) & 
            (bgr_psl['bgr'] == row['bgr']),
            'psl_mrp'
        ]
        if not matching_row.empty:
            osa_tab.loc[index, 'psl_mrp'] = round(matching_row.values[0],0)
    
    osa_tab=osa_tab.merge(d['full_city_df'].groupby(['gc_platform', 'bgr'])['gc_city'].apply(lambda x: ','.join(x.head(2).unique())).reset_index(name='top_geo'),how='left',on=['gc_platform','bgr'])
    # bgr_city_psl=f.groupby(['gc_platform','bgr','gc_city'])['psl_mrp'].sum().reset_index()

    osa_tab['gc_city']=osa_tab['top_geo'].str.split(',').str[0]
    
    for index,row in osa_tab.iterrows():
        osa_tab.loc[index,'cat_share_dr1']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'share_dr1'].values[0]
        osa_tab.loc[index,'cat_change_market_share']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'pct_change'].values[0]
        osa_tab.loc[index,'cat_change_wt_osa']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'change_wt_osa'].values[0]
        osa_tab.loc[index,'cat_wt_osa_dr1']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'wt_osa_dr1'].values[0]
        osa_tab.loc[index,'cat_sales_impact']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'sales_impact'].values[0]

        if len(row['top_geo'].split(','))==1:
            osa_tab.loc[index,'city1sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[0]),'sales_impact'].values[0]
            osa_tab.loc[index,'city2sl']=0
            osa_tab.loc[index,'city1impact']=round(100*osa_tab.loc[index,'city1sl']/osa_tab.loc[index,'cat_sales_impact'],0)
            osa_tab.loc[index,'city2impact']=0
            osa_tab.loc[index,'city1']=row['top_geo'].split(',')[0]
            osa_tab.loc[index,'top_geo']=f"{row['top_geo'].split(',')[0]} ({int(osa_tab.loc[index,'city1impact'])}%)"
        else:
            osa_tab.loc[index,'city1']=row['top_geo'].split(',')[0]
            osa_tab.loc[index,'city2']=row['top_geo'].split(',')[1]
            osa_tab.loc[index,'city1sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[0]),'sales_impact'].values[0]
            osa_tab.loc[index,'city2sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[1]),'sales_impact'].values[0]
            osa_tab.loc[index,'city1impact']=round(100*osa_tab.loc[index,'city1sl']/osa_tab.loc[index,'cat_sales_impact'],0)
            osa_tab.loc[index,'city2impact']=round(100*osa_tab.loc[index,'city2sl']/osa_tab.loc[index,'cat_sales_impact'],0)
            osa_tab.loc[index,'top_geo']=f"{row['top_geo'].split(',')[0]} ({int(osa_tab.loc[index,'city1impact'])}%) , {row['top_geo'].split(',')[1]} ({int(osa_tab.loc[index,'city2impact'])}%)"
        
        # osa_tab.loc[index,'cchange_wt_osa']=d['city_df_brand'].loc[(d['city_df_brand']['gc_platform']==row['gc_platform']) & (d['city_df_brand']['bgr']==row['bgr']) & (d['city_df_brand']['gc_city']==row['gc_city']) ,'change_wt_osa'].values[0]
    
    # print(osa_tab.head())
    osa_tab=osa_tab[osa_tab['change_wt_osa']<0]
    sku_psl_loss=f.groupby(['gc_platform','bgr','gc_city','product_id','title'])['psl_mrp'].sum().reset_index()
    df_sku=pd.DataFrame()
    for index, row in osa_tab.iterrows():
        sku_df=sku_psl_loss.loc[(sku_psl_loss['gc_platform']==row['gc_platform']) & (sku_psl_loss['bgr']==row['bgr']) & (sku_psl_loss['gc_city']==row['gc_city'])]
        if sku_df.empty:
            continue
        sku_df=sku_df.sort_values(by='psl_mrp',ascending=False).head(1)
        osa_tab.loc[index,'product_id']=sku_df['product_id'].values[0]
        # osa_tab.loc[index,'psl_mrp']=sku_df['psl_mrp'].values[0]
        osa_tab.loc[index,'title']=sku_df['title'].values[0]

# st_ids=tuple(f['store_id'].unique())

    for index,row in osa_tab.iterrows():
        pid=row['product_id']
        bgr=row['bgr']
        platform=row['gc_platform']
        city=row['gc_city']
        title=row['title']
        df_sk=sku_data1.loc[(sku_data1['product_id']==pid)  & (sku_data1['gc_platform']==platform) & (sku_data1['gc_city']==city)]
        osa_tab.loc[index,'high_psl_stores']=df_sk['high_psl_stores'].values[0]
        osa_tab.loc[index,'listed_in']=df_sk['listed_in'].values[0]
        # dcdf=pd.concat([dcdf,df_sk],ignore_index=True)

    def get_doi(acc,osa_tab,l1_start,l1_end,l2_start,l2_end):
        dfs=[]
        db_name = acc['db_name']
        for _, row in osa_tab.iterrows():


            platform = row['gc_platform']
            city = row['gc_city']
            bgr = row['bgr']
            product_id = row['product_id']
            title = row['title']
            # get plain YYYY-MM-DD strings
            l1s = l1_start.date().isoformat()
            l1e = l1_end.date().isoformat()
            l2e = l2_end.date().isoformat()
            l2s = l2_start.date().isoformat()
            if platform=='blinkit':
                query=f"""
                        select lower(gc_platform) as gc_platform,lower(bgr) as bgr,gc_city,sku as title,product_id,round(avg(backend_doi_value),0) as bdoi,round(avg(frontend_doi_value),0) as fdoi 
                        from {db_name}.aggregate.{platform}_qcom_facility_doi_metrics 
                        where (snapshot_date between DATE '{l1s}' and DATE'{l1e}') and gc_city='{city}' and regexp_replace(lower(bgr), '[^a-z0-9 ]', '') =$${bgr}$$ and product_id='{product_id}'
                        group by 1,2,3,4,5
            """
            elif platform=='instamart':
                query=f"""
                        select lower(gc_platform) as gc_platform,lower(bgr) as bgr,gc_city,sku as title,product_id,round(avg(backend_doi_value),0) as bdoi,round(avg(frontend_doi_value),0) as fdoi 
                        from {db_name}.aggregate.{platform}_qcom_facility_doi_metrics 
                        where (snapshot_date between DATE '{l1s}' and DATE'{l1e}') and gc_city='{city}' and regexp_replace(lower(bgr), '[^a-z0-9 ]', '') =$${bgr}$$ and product_id='{product_id}'
                        group by 1,2,3,4,5
            """
            elif platform=='zepto':
                query=f"""
                        select lower(gc_platform) as gc_platform,lower(bgr) as bgr,gc_city,sku as title,product_id,round(avg(backend_doi_value),0) as bdoi,round(avg(frontend_doi_value),0) as fdoi 
                        from {db_name}.aggregate.{platform}_qcom_facility_doi_metrics 
                        where (snapshot_date between DATE '{l1s}' and DATE'{l1e}') and gc_city='{city}' and regexp_replace(lower(bgr), '[^a-z0-9 ]', '') =$${bgr}$$ and product_id='{product_id}'
                        group by 1,2,3,4,5
            """
            else:
                raise ValueError(f"Platform {platform} not supported")
            df=sfFetch(query)
            # print(df)
            df.columns=df.columns.str.lower()
            dfs.append(df)
        return pd.concat(dfs,ignore_index=True)
    doi=get_doi(acc,osa_tab,l1_start,l1_end,l2_start,l2_end)
    print(doi.head())
    osa_tab['psl_mrp']=osa_tab['psl_mrp'].apply(lambda x: f"₹{x/1e7:.1f} Cr" if x >= 1e7 else f"₹{x/1e5:.1f} L")

    full_osa=osa_tab.merge(doi,on=['gc_platform','bgr','gc_city','product_id','title'],how='left')

    
    
    def create_osa_recommendations(full_osa):
        df=pd.DataFrame()

        for index,row in full_osa.iterrows():
            df.loc[index,'gc_platform']=row['gc_platform']
            df.loc[index,'bgr']=row['bgr']
            df.loc[index,'top_geo']=row['top_geo']
            df.loc[index,'share_dr1']=f"{round(row['cat_share_dr1'],1)}%"
            df.loc[index,'pct_change']=f"{round(row['cat_change_market_share'],1)}%"
            df.loc[index,'wt_osa_dr1']=f"{round(row['cat_wt_osa_dr1'],1)}%"
            df.loc[index,'change_wt_osa']=f"{round(row['cat_change_wt_osa'],1)}%"
            df.loc[index,'what_happened']=f"OSA dropped from {int(row['wt_osa_dr1']-row['change_wt_osa'])}% to {(row['wt_osa_dr1']):.1f}%\nTop impacted cities are {row['top_geo']}"
            df.loc[index,'potential_sales_loss']=row['psl_mrp']
            rec_text=f"OSA Recommendations:\n- {row['title']}:\n"
            if (row['high_psl_stores']-row['listed_in'])>5:
                rec_text=rec_text+f"• Fix DS listing in {int(row['high_psl_stores']-row['listed_in'])} stores in {row['gc_city']}\n"
            if row['bdoi']==0:
                rec_text=rec_text+ f"• No stock available, fix backend inventory in {row['gc_city']}\n"
            elif row['bdoi']<5:
                rec_text=rec_text+f"• Stock-out risk, Only {row['bdoi']} days of stock left in {row['gc_city']}; backend inventory needs to be refilled\n"

            else:
                if row['fdoi']<3:
                    rec_text=rec_text+f"• Fix transfer issue from warehouse to DS in {row['gc_city']}\n"
                
            
            df.loc[index,'recommendation']=rec_text
            if len(rec_text.splitlines())<3:
                df.loc[index,'recommendation']=''
            

        return df
    osa_rec=create_osa_recommendations(full_osa)
    return osa_rec


def format_osa(final_osa):
    if final_osa.empty:
        return pd.DataFrame([{
            f"2*{i}": "" for i in range(1, 17)
        }])

    row1 = final_osa.iloc[0]
    out = {
        '2*1': f"{row1['bgr']} ({row1['gc_platform']})",
        '2*2': row1['share_dr1'],
        '2*3': row1['pct_change'],
        '2*4': row1['wt_osa_dr1'],
        '2*5': row1['change_wt_osa'],
        '2*6': row1['potential_sales_loss'],
        '2*7': row1['top_geo'],
        '2*8': row1['recommendation'],
    }

    if len(final_osa) > 1:
        row2 = final_osa.iloc[1]
        out.update({
            '2*9':  f"{row2['bgr']} ({row2['gc_platform']})",
            '2*10': row2['share_dr1'],
            '2*11': row2['pct_change'],
            '2*12': row2['wt_osa_dr1'],
            '2*13': row2['change_wt_osa'],
            '2*14': row2['potential_sales_loss'],
            '2*15': row2['top_geo'],
            '2*16': row2['recommendation'],
        })
    else:
        out.update({
            '2*9': "", '2*10': "", '2*11': "", '2*12': "",
            '2*13': "", '2*14': "", '2*15': "", '2*16': ""
        })

    return pd.DataFrame([out])
    

def get_summary(d,acc,l1_start,l1_end,l2_start,l2_end,f1):
    def fmt_amt(x):
        if x >= 1e7:
            return f"₹{x/1e7:.1f} Cr "
        else:
            return f"₹{x/1e5:.1f} L"
    def fmt_change(x):
        if x > 0:
            return f"(↑ {fmt_amt(np.abs(x))})"
        else:
            return f"(↓ {fmt_amt(np.abs(x))})"
    ms=d['ms']
    top2=d['top2']
    below2=d['below2']
    new_stores=sfFetch(f"""WITH store_data AS (
    SELECT 
        gc_platform,
        store_id,
        LOWER(TRIM(gc_city)) AS gc_city_fin,
        MIN(gc_crawl_date) AS store_first_seen
    FROM 
        wholetruth.aggregate.daily_store_catalog
   
    GROUP BY 
        gc_platform, 
        store_id, 
        LOWER(TRIM(gc_city))
    having min(gc_crawl_date) between DATE '{l2_start}' and DATE '{l1_end}'
)
SELECT * 
FROM store_data
ORDER BY gc_platform;""")
    new_stores.columns=new_stores.columns.str.lower()
    new_stores['gc_platform']=new_stores['gc_platform'].str.lower()
    your_stores=sfFetch(f"""select distinct gc_platform,store_id from {acc['db_name']}.aggregate.qcom_locality_metrics where snapshot_date = '{l1_end}' and  is_own_brand=1 and total_listed=1""")
    if your_stores.empty:
        your_stores=sfFetch(f"""select distinct gc_platform,store_id from {acc['db_name']}.aggregate.qcom_locality_metrics where snapshot_date = Date'{l1_end}' -1 and  is_own_brand=1 and total_listed=1""")
    your_stores.columns=your_stores.columns.str.lower()
    ns=new_stores.groupby('gc_platform')['store_id'].nunique().reset_index(name='new_stores')
    for index,plat in ns.iterrows():
        platform=plat['gc_platform']
        your_stores_df=your_stores.loc[your_stores['gc_platform']==platform]
        new_stores_df=new_stores.loc[new_stores['gc_platform']==platform]
        v=new_stores_df[new_stores_df['store_id'].isin(your_stores_df['store_id'].unique())]['store_id'].nunique()
        # v=your_stores_df[your_stores_df['store_id'].isin(new_stores_df['store_id'].unique())]['store_id'].nunique()
        ns.at[index,'listed_in_new_stores']=v

    def get_summary_vars(ms):
        """
        Compute summary metrics for each platform from the `ms` DataFrame.

        Returns a DataFrame with:
        - gc_platform
        - Offtake: current-period sum (dr2)
        - prev_Offtake: prior-period sum (dr1)
        - share_delta_pct: percent change arrow between dr1 and dr2
        """
        # Filter to own-brand only
        df = ms[ms['is_own_brand'] == 1].copy()

        # Period labels
        dr1, dr2 = 'dr1', 'dr2'

        # Total off-take sums
        tot_dr1 = df.loc[df['date_period'] == dr1, 'actual_sum'].sum()
        tot_dr2 = df.loc[df['date_period'] == dr2, 'actual_sum'].sum()

        # Compute percent change and arrow indicator
        change = tot_dr1 - tot_dr2 

        arrow = '↑' if change > 0 else '↓'
        share_delta_str = f"{arrow} {abs(change)}%"

        # Per-platform sums
        pt = (
            df
            .groupby(['gc_platform', 'date_period'])['actual_sum']
            .sum()
            .unstack(fill_value=0)
            .reset_index()
        )
        pt = pt.rename(columns={dr1: 'curr_sum', dr2: 'prev_sum'})
        def fmt_amt(x):
            if x >= 1e7:
                return f"₹{x/1e7:.1f} Cr "
            else:
                return f"₹{x/1e5:.1f} L"
        def fmt_change(x):
            if x > 0:
                return f"(↑ {fmt_amt(np.abs(x))})"
            else:
                return f"(↓ {fmt_amt(np.abs(x))})"
        
        
        pt['share_delta'] = pt['curr_sum'] - pt['prev_sum']
        # pt['curr_sum'] = pt['curr_sum'].apply(fmt_amt)
        pt['prev_sum'] = pt['prev_sum'].apply(fmt_amt)
        pt['share_delta_str']=pt['share_delta'].apply(fmt_change)

        # Build result
        summary = pd.DataFrame({
            'gc_platform': pt['gc_platform'],
            'Offtake': pt['curr_sum'],
            'prev_Offtake': pt['prev_sum'],
            'share_delta': pt['share_delta_str']
        })
        return summary


    def compile_offtake_summary(ms, below2,top2, ns):
        """
        Combine platform off-take metrics, top impacted categories, and new DS coverage into a summary table.

        Args:
        - ms: DataFrame with columns including date_period, is_own_brand, actual_sum
        - below2: DataFrame with ['gc_platform','bgr'] of impacted categories
        - ns: DataFrame with ['gc_platform','new_stores','listed_in_new_stores']

        Returns:
        A DataFrame with:
            gc_platform,
            Offtake (formatted ₹),
            Top impacted categories,
            New DS coverage (stores and listing %)
        """
        # Base metrics
        base = get_summary_vars(ms)

        # base['Offtake']=base['Offtake']+'\n'+base['share_delta']
# 1. Build your formatted labels
        top2['formatted'] = top2.apply(lambda x: f"{x['bgr']}({x['pct_change']:.1f}%)", axis=1)
        below2['formatted'] = below2.apply(lambda x: f"{x['bgr']}({x['pct_change']:.1f}%)", axis=1)

        # 2. Aggregate each to one row per platform
        gainers = (
            top2
            .groupby('gc_platform')
            .agg({
                'bgr'         :'first',
                'pct_change'  :'first',
                'formatted'   :'first'
            })
            .reset_index()
            .rename(columns={'formatted':'gainers'})
        )

        drainers = (
            below2
            .groupby('gc_platform')
            .agg({
                'bgr'         :'first',
                'pct_change'  :'first',
                'formatted'   :'first'
            })
            .reset_index()
            .rename(columns={'formatted':'drainers'})
        )

        # 3. Merge everything so you always get both columns
        merged = pd.merge(
            gainers,
            drainers,
            on='gc_platform',
            how='outer',
            suffixes=('','_dr')
        )

        # 4. Fill any holes so we can operate safely
        merged[['bgr','bgr_dr']]                 = merged[['bgr','bgr_dr']].fillna('')
        merged[['pct_change','pct_change_dr']]   = merged[['pct_change','pct_change_dr']].fillna(0)
        merged[['gainers','drainers']]           = merged[['gainers','drainers']].fillna('')

        # 5. If it’s the *same* brand on both sides, keep only the “winning” formatted entry…
        same = merged['bgr'] == merged['bgr_dr']

        # …if pct_change>0, it’s a gainer
        merged.loc[same & (merged['pct_change'] > 0), 'drainers'] = ''
        # …otherwise it’s a drainer
        merged.loc[same & (merged['pct_change'] <= 0), 'gainers']   = ''

        # 6. Now pull out exactly the columns you want
        cat = merged[['gc_platform','gainers','drainers']]

        # # Use vectorized operations with mask
        # mask = (gainers['gainers'] == drainers['drainers']) & (gainers['pct_change'] > 0)
        # gainers.loc[mask, 'gainers'] = gainers['gainers'] + gainers['pct_change'].apply(lambda x: f"(↑{x:.1f}%)")
        # drainers.loc[mask, 'drainers'] = ''

        # mask = (gainers['gainers'] == drainers['drainers']) & (gainers['pct_change'] <= 0)
        # drainers.loc[mask, 'drainers'] = drainers['drainers'] + drainers['pct_change'].apply(lambda x: f"(↓{abs(x):.1f}%)")
        # gainers.loc[mask, 'gainers'] = ''



        # cat=pd.merge(gainers.drop(columns='pct_change'), drainers.drop(columns='pct_change'), on='gc_platform', how='outer')
        # # cat['Top impacted categories'] = cat.apply(format_top_impacted, axis=1)


# Combine into single impacted_categories column
        

        # cat = (
        #     below2
        #     .groupby('gc_platform')['bgr']
        #     .unique()
        #     .apply(lambda arr: ", ".join(arr))
        #     .reset_index(name='Top impacted categories')
        # )
        ns['new_stores']=ns['new_stores'].astype(int)
        ns['listing_pct'] = ns.apply(
            lambda r: f"{int(round(100 * r['listed_in_new_stores'] / r['new_stores'], 0))}%",
            axis=1
        )
        ns['new_stores']=ns.apply(lambda r: f"out of {r['new_stores']} stores",axis=1)

        ns_copy = ns.copy()
        
        # ns_copy['New DS coverage'] = ns_copy.apply(
        #     lambda r: f"{int(r['new_stores'])} stores ({100*r['listed_in_new_stores']/r['new_stores']:.0f}% listing)",
        #     axis=1
        # )
        ns_out = ns_copy[['gc_platform', 'new_stores','listing_pct']]

        final = (
            base
            .merge(cat, on='gc_platform', how='left')
            .merge(ns_out, on='gc_platform', how='left')
        )

        return final[['gc_platform', 'Offtake','share_delta', 'gainers', 'drainers', 'new_stores','listing_pct']]
    df_summary=compile_offtake_summary(ms, below2,top2, ns)

    df_summary=df_summary.merge(f1,on='gc_platform',how='left')
    df_summary['psl_per']=round(100*df_summary['psl_mrp']/df_summary['Offtake'],1)
    df_summary['Offtake']=df_summary['Offtake'].apply(fmt_amt)
    df_summary['psl_mrp']=df_summary['psl_mrp'].apply(lambda x: f"₹{x/1e7:.1f} Cr" if x >= 1e7 else f"₹{x/1e5:.1f} L")
    df_summary['psl_mrp']=df_summary.apply(lambda x: f"{x['psl_mrp']}\n({x['psl_per']}% of offtake)",axis=1)
    df_summary=df_summary.drop(columns='date_period')
    df_summary=df_summary[['gc_platform','Offtake','share_delta','gainers','drainers','psl_mrp','listing_pct','new_stores']]
    # 1) Define your platforms & desired order
    platforms = ['blinkit', 'instamart', 'zepto']

    # 2) Reindex df_summary to guarantee one row per platform
    #    any missing platform will become a row of NaNs
    df_fixed = (
        df_summary
        .set_index('gc_platform')        # index by platform
        .reindex(platforms)              # enforce the 3 rows
        .reset_index()                   # bring platform back as column
    )

    # 3) Fill any missing data with blank strings
    df_fixed = df_fixed.fillna('')

    # 4) Drop the platform column (you flatten it out of the markers)
    df = df_fixed.drop(columns=['gc_platform'])

    # 5) Now flatten row-major to a single row of length 3 * n_cols
    values = df.to_numpy().flatten(order='C')   # length = 3 * number_of_columns

    # 6) Build your marker names 1*1 … 1*21
    col_names = [f"1*{i}" for i in range(1, len(values) + 1)]

    single_df_summary = pd.DataFrame([values], columns=col_names)




    
    return single_df_summary

def get_relevant_keywords(acc, l1_start, l1_end, l2_start, l2_end):

    # db_name = acc['db_name']
  db_name=acc['db_name']
  # ISO date strings
  l1s = l1_start.date().isoformat()
  l1e = l1_end.date().isoformat()
  l2s = l2_start.date().isoformat()
  l2e = l2_end.date().isoformat()

  # Fetch both LMTD and MTD periods
  query = f"""
WITH ranked_keywords AS (
SELECT
  CASE
    WHEN snapshot_date BETWEEN DATE '{l2s}' AND DATE '{l2e}' THEN 'dr2'
    WHEN snapshot_date BETWEEN DATE '{l1s}' AND DATE '{l1e}' THEN 'dr1'
  END AS date_period,
  gc_platform,
  bgr,
  gc_city,
  keyword,
  SUM(is_own_brand * ad_impressions)     AS own_imp,
  SUM(total_impressions)                 AS tot_imp,
  SUM(ad_impressions)                    AS ad_imp,
  CASE WHEN SUM(ad_impressions) = 0 THEN 0
        ELSE ROUND(SUM(is_own_brand * ad_impressions) * 100.0 / SUM(ad_impressions), 1)
  END                                    AS ad_sov,
  MAX(keyword_weight)                    AS max_keyword_weight
FROM {db_name}.aggregate.qcom_brand_keyword_impressions
WHERE
  snapshot_date BETWEEN DATE '{l2s}' AND DATE '{l2e}'
  OR snapshot_date BETWEEN DATE '{l1s}' AND DATE '{l1e}'
GROUP BY 1,2,3,4,5
),

metrics AS (
SELECT
  date_period,
  gc_platform,
  bgr,
  gc_city,
  keyword,
  max_keyword_weight,
  ad_sov,
  ROUND(100.0 * tot_imp 
    / NULLIF(SUM(tot_imp) OVER (PARTITION BY date_period, gc_platform, bgr, gc_city), 0)
  , 1) AS cat_share,
  ROUND(100.0 * max_keyword_weight 
    / NULLIF(SUM(max_keyword_weight) OVER (PARTITION BY date_period, gc_platform, bgr, gc_city), 0)
  , 1) AS vol_share
FROM ranked_keywords
)
SELECT *
FROM metrics;
"""

  # Raw DataFrame for both periods
  df_raw = sfFetch(query)
  df_raw.columns = df_raw.columns.str.lower()

  # Pivot ad_sov to get last month and this month SOV
  pivot = (
      df_raw
      .pivot_table(
          index=['gc_platform', 'bgr', 'gc_city', 'keyword'],
          columns='date_period',
          values='ad_sov'
      )
      .reset_index()
      .rename(columns={'dr2': 'sov_dr2', 'dr1': 'sov_dr1'})
  )
  pivot.columns.name = None

  # Focus on MTD rows and compute relevance & final scores
  df_mtd = df_raw[df_raw['date_period'] == 'dr1'].copy()
  df_mtd['relevant_score'] = np.where(
      df_mtd['vol_share'] != 0,
      np.round(df_mtd['cat_share'] / df_mtd['vol_share'], 2),
      0
  )
  df_mtd_non_relevant = df_mtd[(df_mtd['relevant_score']<=1)&(df_mtd['relevant_score']>0) & (df_mtd['vol_share']>1)]


  df_mtd = df_mtd[(df_mtd['relevant_score'] > 1) & (df_mtd['vol_share']>1)]
  df_mtd['final_score'] = df_mtd['vol_share'] * np.log1p(df_mtd['relevant_score'])

  # Sort by final_score descending within each group
  df_mtd = df_mtd.sort_values(['gc_platform', 'bgr', 'gc_city', 'final_score'],ascending=[True, True, True, False])
  df_mtd_non_relevant['final_score'] = df_mtd_non_relevant['vol_share'] * np.log1p(df_mtd_non_relevant['relevant_score'])
  df_mtd_non_relevant = df_mtd_non_relevant.sort_values(
      ['gc_platform', 'bgr', 'gc_city', 'final_score'],
      ascending=[True, True, True, False]
  )
  df_non_relevant = df_mtd_non_relevant.merge(
      pivot,
      on=['gc_platform', 'bgr', 'gc_city', 'keyword'],
      how='left'
  )
  # Merge to bring in last month's SOV
  df = df_mtd.merge(
      pivot,
      on=['gc_platform', 'bgr', 'gc_city', 'keyword'],
      how='left'
  )

  # Build recommendations per gc_platform, bgr, gc_city
  recs = []
  for (platform, bgr, city), grp in df.groupby(['gc_platform', 'bgr', 'gc_city']):
      # Keywords where share of voice dropped
      drop_kws = grp.loc[grp['sov_dr1'] < grp['sov_dr2'], 'keyword'].tolist()
      if drop_kws:
          kws_to_bid = drop_kws[:2]
      else:
          # Top candidate by final_score
          kws_to_bid = grp['keyword'].head(1).tolist()
      recs.append({
          'gc_platform': platform,
          'bgr': bgr,
          'gc_city': city,
          'recommended_keywords': kws_to_bid
      })

  recommendations = pd.DataFrame(recs)
  return df, recommendations,df_non_relevant

def get_kw_brand_sov(acc,l1_start,l1_end,l2_start,l2_end):
  kw_brand_sov=sfFetch(f"""

  with t as (SELECT
    CASE
      WHEN snapshot_date BETWEEN DATE '{l1_start}' AND DATE '{l1_end}' THEN 'dr1'
      WHEN snapshot_date BETWEEN DATE '{l2_start}' and DATE '{l2_end}' THEN 'dr2'
    END AS date_period,
    gc_platform,
    bgr,
    gc_city,
    keyword,
    brand,
    SUM(ad_impressions)     AS own_imp,
    
  FROM {acc['db_name']}.aggregate.qcom_brand_keyword_impressions
  WHERE
    snapshot_date BETWEEN DATE '{l2_start}' AND DATE '{l2_end}'
    OR snapshot_date BETWEEN DATE '{l1_start}' AND DATE '{l1_end}'
  GROUP BY 1,2,3,4,5,6)
  ,
  s as (
  select 
  CASE
      WHEN snapshot_date BETWEEN DATE '{l2_start}' AND DATE '{l2_end}' THEN 'dr2'
      WHEN snapshot_date BETWEEN DATE '{l1_start}' AND DATE '{l1_end}' THEN 'dr1'
    END AS date_period,
  gc_platform,bgr,gc_city,keyword,sum(ad_impressions) tot_imp from {acc['db_name']}.aggregate.qcom_brand_keyword_impressions
  WHERE
    snapshot_date BETWEEN DATE '{l2_start}' AND DATE '{l2_end}'
    OR snapshot_date BETWEEN DATE '{l1_start}' AND DATE '{l1_end}' 
    group by 1,2,3,4,5
  ),
  z as(select t.date_period,
  t.gc_platform,t.bgr,t.gc_city,t.keyword,t.brand ,round(100*t.own_imp/nullif(tot_imp,0),1) ad_sov
  from t left join s  USING(date_period,gc_platform,bgr,gc_city,keyword))
  ,
  brand_info as (select distinct  gc_platform,brand ,is_own_brand from {acc['db_name']}.aggregate.qcom_product_dimensions)
  select z.* ,bi.is_own_brand from z left join brand_info bi  USING (gc_platform,brand) """)
  kw_brand_sov.columns=kw_brand_sov.columns.str.lower()
  kw_brand_sov['gc_platform']=kw_brand_sov['gc_platform'].str.lower()
  kw_brand_sov=kw_brand_sov.fillna(0)
  return kw_brand_sov
# a,b,c=get_relevant_keywords(acc,l1_start,l1_end,l2_start,l2_end)
# kw_brand_sov=get_kw_brand_sov(acc,l1_start,l1_end,l2_start,l2_end)


def get_contribution(acc,l1_start,l1_end,l2_start,l2_end):
    df=sfFetch(f"""
    WITH pid_totals AS (
  SELECT
  product_id,
  title,
    gc_platform,
    bgr,
    gc_city,
    SUM(reported_offtake_mrp) AS total_mrp
  FROM {acc['db_name']}.aggregate.qcom_pid_metrics
  WHERE
    snapshot_date BETWEEN DATE '{l1_start}' AND DATE '{l1_end}'
    AND is_own_brand = 1
  GROUP BY
    product_id,
    title,
    gc_platform,
    bgr,
    gc_city
)
SELECT
  product_id,
  title,
  gc_platform,
  bgr,
  gc_city,
  ROUND(
    100.0 * total_mrp
          / NULLIF(
              SUM(total_mrp) OVER (PARTITION BY gc_platform, bgr, gc_city),
              0
            ),
    1
  ) AS contribution_pct
FROM pid_totals
ORDER BY
gc_platform,
  bgr,
  gc_city,
  
  contribution_pct DESC
;
    """)
    return df


def get_sov_summary(d,acc,l1_start,l1_end,l2_start,l2_end):
    bgr_data=get_bgr_data(acc,l1_start,l1_end,l2_start,l2_end)
    own_bgr_data=bgr_data[bgr_data['is_own_brand']==1]
    a,b,c=get_relevant_keywords(acc,l1_start,l1_end,l2_start,l2_end)
    kw_brand_sov=get_kw_brand_sov(acc,l1_start,l1_end,l2_start,l2_end)
    below2=d['below2'].sort_values(by='sales_impact',ascending=True).merge(own_bgr_data,on=['gc_platform','bgr','brand','is_own_brand'],how='left')
    sov_tab=d['city_df_brand'].sort_values(by='sales_impact',ascending=True).query('change_sov<0')

    sov_tab=sov_tab.merge(d['full_city_df'].groupby(['gc_platform', 'bgr'])['gc_city'].apply(lambda x: ','.join(x.head(2).unique())).reset_index(name='top_geo'),how='left',on=['gc_platform','bgr'])
    # sov_tab['gc_city']=sov_tab['top_geo'].str.split(',').str[0]
    for index,row in sov_tab.iterrows():
        sov_tab.loc[index,'cat_share_dr1']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'share_dr1'].values[0]
        sov_tab.loc[index,'cat_change_market_share']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'pct_change'].values[0]
        sov_tab.loc[index,'cat_change_sov']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'change_sov'].values[0]
        sov_tab.loc[index,'cat_sov']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'sov_dr1'].values[0]
        sov_tab.loc[index,'cat_sales_impact']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'sales_impact'].values[0]

        if len(row['top_geo'].split(','))==1:
            sov_tab.loc[index,'city1sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[0]),'sales_impact'].values[0]
            sov_tab.loc[index,'city2sl']=0
            sov_tab.loc[index,'city1impact']=round(100*sov_tab.loc[index,'city1sl']/sov_tab.loc[index,'cat_sales_impact'],0)
            sov_tab.loc[index,'city2impact']=0
            sov_tab.loc[index,'city1']=row['top_geo'].split(',')[0]
            sov_tab.loc[index,'top_geo']=row['top_geo']
            sov_tab.loc[index,'top_geo_str']=f"{row['top_geo'].split(',')[0]}"
        else:
            sov_tab.loc[index,'city1']=row['top_geo'].split(',')[0]
            sov_tab.loc[index,'city2']=row['top_geo'].split(',')[1]
            sov_tab.loc[index,'city1sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[0]),'sales_impact'].values[0]
            sov_tab.loc[index,'city2sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[1]),'sales_impact'].values[0]
            sov_tab.loc[index,'city1impact']=round(100*sov_tab.loc[index,'city1sl']/sov_tab.loc[index,'cat_sales_impact'],0)
            sov_tab.loc[index,'city2impact']=round(100*sov_tab.loc[index,'city2sl']/sov_tab.loc[index,'cat_sales_impact'],0)
            sov_tab.loc[index,'top_geo']=row['top_geo']
            sov_tab.loc[index,'top_geo_str']=f"{row['top_geo'].split(',')[0]}"

        
    cnt_df=get_contribution(acc,l1_start,l1_end,l2_start,l2_end)
    cnt_df.columns=cnt_df.columns.str.lower()
    cnt_df['gc_platform']=cnt_df['gc_platform'].str.lower()
    cnt_df['bgr']=cnt_df['bgr'].str.lower()
    sov_tab=sov_tab[sov_tab['change_sov']<0]
    # print(sov_tab.to_markdown())
    if sov_tab.empty:
        print('City level SOV drop not found')
        return pd.DataFrame(columns=['gc_platform', 'bgr', 'What happened?', 'RCA', 'Recommendation','share_dr1', 'pct_change', 'sov_dr1', 'change_sov'])
    ppu_cat=get_ppu_cat_brand(acc,l1_start,l1_end,l2_start,l2_end)
    ppu_cat_comp=ppu_cat[ppu_cat['is_own_brand']==0]
    # sov_tab=sov_tab.merge(ppu_cat,on=['account','gc_platform','bgr','brand','is_own_brand'],how='left')
    b['gc_platform']=b['gc_platform'].str.lower()
    b['bgr']=b['bgr'].str.lower()
    c['gc_platform']=c['gc_platform'].str.lower()
    c['bgr']=c['bgr'].str.lower()
    for index, row in sov_tab.iterrows():
        sku_df=d['a'].query('is_own_brand==1').loc[(d['a'].query('is_own_brand==1')['gc_platform']==row['gc_platform']) & (d['a'].query('is_own_brand==1')['bgr']==row['bgr']) & (d['a'].query('is_own_brand==1')['gc_city']==row['gc_city'])].merge(cnt_df,on=['gc_platform','bgr','gc_city','product_id','title'],how='left')
        sku_df=sku_df[sku_df['contribution_pct']>5].query('change_sov<0 and change_wt_osa<0').sort_values(by='contribution_pct',ascending=False).head(1)
        if sku_df.empty:
            continue
        sov_tab.loc[index,'product_id']=sku_df['product_id'].values[0]
        # sov_tab.loc[index,'psl_mrp']=sku_df['psl_mrp'].values[0]
        sov_tab.loc[index,'title']=sku_df['title'].values[0]
        sov_tab.loc[index,'sku_change_sov']=sku_df['change_sov'].values[0]
        sov_tab.loc[index,'sku_change_wt_osa']=sku_df['change_wt_osa'].values[0]
    for index, row in sov_tab.iterrows():
        comps = ppu_cat_comp.loc[
            (ppu_cat_comp['gc_platform'] == row['gc_platform']) & 
            (ppu_cat_comp['bgr'] == row['bgr']) & 
            (ppu_cat_comp['ppu_category'] == row['ppu_category']), 
            'brand'
        ].unique()
        
        target_osa=d['merged_comp'].loc[
                (d['merged_comp']['gc_platform'] == row['gc_platform']) & 
                (d['merged_comp']['bgr'] == row['bgr']) & (d['merged_comp']['gc_city']==row['gc_city']) &
                (d['merged_comp']['brand'].isin(comps))]['wt_osa_dr1'].max()
        if target_osa is np.nan:
            comps = ppu_cat_comp.loc[
            (ppu_cat_comp['gc_platform'] == row['gc_platform']) & 
            (ppu_cat_comp['bgr'] == row['bgr']), 
            'brand'
        ].unique()
            target_osa=d['merged_comp'].loc[
                (d['merged_comp']['gc_platform'] == row['gc_platform']) & 
                (d['merged_comp']['bgr'] == row['bgr']) & (d['merged_comp']['gc_city']==row['gc_city']) &
                (d['merged_comp']['brand'].isin(comps))]['wt_osa_dr1'].max()
        if len(comps) > 1:
            
            
            comp_ms = d['merged_comp'].loc[
                (d['merged_comp']['gc_platform'] == row['gc_platform']) & 
                (d['merged_comp']['bgr'] == row['bgr']) & 
                (d['merged_comp']['brand'].isin(comps)) & 
                (d['merged_comp']['change_adsov'] > 0)
            ].sort_values(by='share_dr1', ascending=False).head(4).sort_values(by='change_adsov', ascending=False).head(1)
            if comp_ms.empty:
                comp_ms=None
            else:
                comp_ms=comp_ms['brand'].values[0]
        elif len(comps) == 0:
            continue  # Skip if no competitors found
        else:
            comp_ms = comps[0]  # If only one competitor, use that

        
        # Assigning comp_ms to the DataFrame
        if row['wt_osa_dr1']>target_osa:
            target_osa=row['wt_osa_dr1']+row['change_wt_osa']

        sov_tab.loc[index, 'comp_ms'] = comp_ms
        sov_tab.loc[index, 'target_osa'] = target_osa
        
        # Get keywords for the selected competitor
        kw_data = (
            compare_monthly_metrics(kw_brand_sov, group_cols=['gc_platform', 'bgr', 'brand', 'keyword', 'gc_city'], metric_cols=['ad_sov'])
            .query('change_ad_sov > 0 and brand == @comp_ms')
            .merge(a, on=['gc_platform', 'bgr', 'keyword', 'gc_city'], how='left')
            .sort_values(by='final_score', ascending=False)
            .head(2)['keyword']
            .unique()
        )
        
        # Convert kw_data to a string representation
        kw_data_str = ', '.join(kw_data) if len(kw_data) > 0 else None
        
        # Assigning kw_data_str to the DataFrame
        sov_tab.loc[index, 'comp_keywords'] = kw_data_str
        rec_kw =b.loc[(b['gc_platform']==row['gc_platform']) & (b['bgr']==row['bgr'])  & (b['gc_city']==row['gc_city']),'recommended_keywords'].values[0]
        rec_kw_str = ', '.join(rec_kw) if len(rec_kw) > 0 else None
        sov_tab.loc[index, 'rec_keywords'] = rec_kw_str
        c1=c[(c['ad_sov']>0) & (c['relevant_score']<0.5)]
        drop_kw=c1.loc[(c1['gc_platform']==row['gc_platform']) & (c1['bgr']==row['bgr'])  & (c1['gc_city']==row['gc_city'])].sort_values(by='ad_sov',ascending=False).head(10).sort_values(by='final_score',ascending=True).head(2)['keyword'].values
        drop_kw_str = ', '.join(drop_kw) if len(drop_kw) > 0 else None
        sov_tab.loc[index, 'drop_keywords'] = drop_kw_str

    def summarize_sov_tab(sov_tab: pd.DataFrame) -> pd.DataFrame:
  
        rows = []
        for _, r in sov_tab.iterrows():
            
            if 'change_sov' in r and pd.notna(r['change_sov']):
                before = r['sov_dr1'] - r['change_sov']
                after  = r['sov_dr1']
                what1 = f"• Overall SOV dropped from {before:.1f}% to {after:.1f}%"
            else:
                what1 = f"• Overall SOV: {r['sov_dr1']:.1f}%"
            geos = str(r.get('top_geo', '')).split(',')
            top2 = [g.strip() for g in geos if g.strip()][:2]

            what2 = f"• Top impacted geographies – {r['top_geo_str']}" if top2 else ""
            what_happened = "\n".join([what1, what2]).strip()

            # --- RCA ---
            rca_lines = []
            # 1) OSA drop on key SKU
            pid = r.get('title', '')
            if pid:
                rca_lines.append(f"• {pid} SOV declined {abs(r['sku_change_sov']):.1f}% due to OSA drop" if pid is not np.nan else "")
            # 2) competitor keywords
            if r.get('comp_keywords', '') is None or r.get('comp_keywords', '') == 'None':
                rca_lines.append(f"\n")
            else:
                comps = str(r.get('comp_keywords', '')).split(',')
                comps = [c.strip() for c in comps if c.strip()][:2]
                if comps:
                    if len(comps) == 1:
                        rca_lines.append(f"• {r['comp_ms']} gained SOV on KWs – “{comps[0]}”")
                # else:
            rca_lines.append(what2)   #     rca_lines.append(f"Comp. brand gained SOV on KWs – “{comps[0]}” and “{comps[1]}”")
            rca = "\n".join(rca_lines)

            # --- Recommendation ---
            rec_lines = ["Sov Recommendation:"]
            # 1) restore OSA
            if pid and 'wt_osa_dr1' in r:
                rec_lines.append(f"• {pid}: Restore OSA to {np.ceil(r['target_osa'] / 5) * 5}% to improve organic SOV" if (pid is not np.nan)  and (r['sku_change_wt_osa']<=-5) else "")
            # 2) increase bids
            recs = str(r.get('rec_keywords', '')).split(',')
            recs = [c.strip() for c in recs if c.strip()][:2]
            if recs:
                if len(recs) == 1:
                    rec_lines.append(f"• Increase Ad SOV on “{recs[0]}” KWs")
                else:
                    rec_lines.append(f"• Increase Ad SOV on “{recs[0]}” and “{recs[1]}” KWs")
            # 3) remove ad spends
            drops = str(r.get('drop_keywords', '')).split(',')
            drops = [d.strip() for d in drops if d.strip()][:1]
            if drops[0] != 'None':
                rec_lines.append(f"• Remove ad spends from {drops[0]}")
            recommendation = "\n".join(rec_lines)
            if len(rec_lines)<2:
                recommendation=''

            # --- Build the summary row ---
            rows.append({
                "gc_platform":           r.get('gc_platform', ""),
                "bgr":            r.get('bgr', ""),
                "What happened?": what_happened,
                "RCA":            rca,
                'rec_lines':rec_lines,
                "Recommendation": recommendation,
                "share_dr1":f'{r.get("cat_share_dr1", "")}%',
                "pct_change":f'{round(r.get("cat_change_market_share", ""),1)}%',
                "sov_dr1":f'{round(r.get("cat_sov", ""),1)}%',
                "change_sov":f'{round(r.get("cat_change_sov", ""),1)}%'
                # "For details":    details
            })

        return pd.DataFrame(rows)
    return summarize_sov_tab(sov_tab)


def format_sov(final_sov):
    if final_sov.empty:
        return pd.DataFrame([{
            f"3*{i}": "" for i in range(1, 15)
        }])

    row1 = final_sov.iloc[0]
    out = {
        '3*1': f"{row1['bgr']} ({row1['gc_platform']})",
        '3*2': row1['share_dr1'],
        '3*3': row1['pct_change'],
        '3*4': row1['sov_dr1'],
        '3*5': row1['change_sov'],
        '3*6': row1['RCA'],
        '3*7': row1['Recommendation']    }

    if len(final_sov) > 1:
        row2 = final_sov.iloc[1]
        out.update({
            '3*8':  f"{row2['bgr']} ({row2['gc_platform']})",
            '3*9': row2['share_dr1'],
            '3*10': row2['pct_change'],
            '3*11': row2['sov_dr1'],
            '3*12': row2['change_sov'],
            '3*13': row2['RCA'],
            '3*14': row2['Recommendation']
            
        })
    else:
        out.update({
            '3*8': "", '3*9': "",'3*10': "", '3*11': "", '3*12': "",
            '3*13': "", '3*14': ""
        })

    return pd.DataFrame([out])
    
def format_disc(final_dis):
    if final_dis.empty:
        return pd.DataFrame([{
            f"4*{i}": "" for i in range(1, 15)
        }])

    row1 = final_dis.iloc[0]
    out = {
        '4*1': f"{row1['bgr']} ({row1['gc_platform']})",
        '4*2': row1['share_dr1'],
        '4*3': row1['pct_change'],
        '4*4': row1['wt_discounting_dr1'],
        '4*5': row1['change_wt_discounting'],
        '4*6': row1['What happened?'],
        '4*7': row1['Why it matters?']    }

    if len(final_dis) > 1:
        row2 = final_dis.iloc[1]
        out.update({
            '4*8':  f"{row2['bgr']} ({row2['gc_platform']})",
            '4*9': row2['share_dr1'],
            '4*10': row2['pct_change'],
            '4*11': row2['wt_discounting_dr1'],
            '4*12': row2['change_wt_discounting'],
            '4*13': row2['What happened?'],
            '4*14': row2['Why it matters?']
            
        })
    else:
        out.update({
            '4*8': "", '4*9': "",'4*10': "", '4*11': "", '4*12': "",
            '4*13': "", '4*14': ""
        })

    return pd.DataFrame([out])
    


def summarize_discount_table(acc,l1_start,l1_end,l2_start,l2_end,d):
    bgr_data=get_bgr_data(acc,l1_start,l1_end,l2_start,l2_end)
    
    own_bgr_data=bgr_data[bgr_data['is_own_brand']==1]
    below2=d['below2'].sort_values(by='sales_impact',ascending=True).merge(own_bgr_data,on=['gc_platform','bgr','brand','is_own_brand'],how='left')
    dis_tab=d['city_df_brand'].sort_values(by='sales_impact',ascending=True).query('change_wt_discounting<0')
    # dis_tab=d['below2'].sort_values(by='sales_impact',ascending=True).merge(own_bgr_data,on=['gc_platform','bgr','brand','is_own_brand'],how='left').query('change_wt_discounting<0')
    

    ppu_cat=get_ppu_cat_brand(acc,l1_start,l1_end,l2_start,l2_end)
    ppu_cat_comp=ppu_cat[ppu_cat['is_own_brand']==0]

    dis_tab=dis_tab.merge(d['full_city_df'].groupby(['gc_platform', 'bgr'])['gc_city'].apply(lambda x: ','.join(x.head(2).unique())).reset_index(name='top_geo'),how='left',on=['gc_platform','bgr'])
    # bgr_city_psl=f.groupby(['gc_platform','bgr','gc_city'])['psl_mrp'].sum().reset_index()

    # dis_tab['gc_city']=dis_tab['top_geo'].str.split(',').str[0]
    
    for index,row in dis_tab.iterrows():
        dis_tab.loc[index,'cat_share_dr1']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'share_dr1'].values[0]
        dis_tab.loc[index,'cat_change_market_share']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'pct_change'].values[0]
        dis_tab.loc[index,'cat_change_wt_discounting']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']) ,'change_wt_discounting'].values[0]
        dis_tab.loc[index,'cat_wt_discounting_dr1']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'wt_discounting_dr1'].values[0]
        dis_tab.loc[index,'cat_sales_impact']=below2.loc[(below2['gc_platform']==row['gc_platform']) & (below2['bgr']==row['bgr']) & (below2['brand']==row['brand']),'sales_impact'].values[0]
        if len(row['top_geo'].split(','))==1:
            dis_tab.loc[index,'city1sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[0]),'sales_impact'].values[0]
            dis_tab.loc[index,'city2sl']=0
            dis_tab.loc[index,'city1impact']=round(100*dis_tab.loc[index,'city1sl']/dis_tab.loc[index,'cat_sales_impact'],0)
            dis_tab.loc[index,'city2impact']=0
            dis_tab.loc[index,'city1']=row['top_geo'].split(',')[0]
            dis_tab.loc[index,'top_geo']=f"{row['top_geo'].split(',')[0]} ({int(dis_tab.loc[index,'city1impact'])}%)"
        else:
            dis_tab.loc[index,'city1']=row['top_geo'].split(',')[0]
            dis_tab.loc[index,'city2']=row['top_geo'].split(',')[1]
            dis_tab.loc[index,'city1sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[0]),'sales_impact'].values[0]
            dis_tab.loc[index,'city2sl']=d['full_city_df'].loc[(d['full_city_df']['gc_platform']==row['gc_platform']) & (d['full_city_df']['bgr']==row['bgr']) & (d['full_city_df']['gc_city']==row['top_geo'].split(',')[1]),'sales_impact'].values[0]
            dis_tab.loc[index,'city1impact']=round(100*dis_tab.loc[index,'city1sl']/dis_tab.loc[index,'cat_sales_impact'],0)
            dis_tab.loc[index,'city2impact']=round(100*dis_tab.loc[index,'city2sl']/dis_tab.loc[index,'cat_sales_impact'],0)
            dis_tab.loc[index,'top_geo']=f"{row['top_geo'].split(',')[0]} ({int(dis_tab.loc[index,'city1impact'])}%) , {row['top_geo'].split(',')[1]} ({int(dis_tab.loc[index,'city2impact'])}%)"
        
    # print(dis_tab.head())
    # dis_tab=dis_tab.merge(ppu_cat,on=['account','gc_platform','bgr','brand','is_own_brand'],how='left')
    # for index,row in dis_tab.iterrows():
        
        # dis_tab.loc[index,'market_share_dr1']=bgr_data.loc[(bgr_data['gc_platform']==row['gc_platform']) & (bgr_data['bgr']==row['bgr']) & (bgr_data['brand']==row['brand']),'market_share_dr1'].values[0]
        # dis_tab.loc[index,'change_market_share']=bgr_data.loc[(bgr_data['gc_platform']==row['gc_platform']) & (bgr_data['bgr']==row['bgr']) & (bgr_data['brand']==row['brand']),'change_market_share'].values[0]
        # dis_tab.loc[index,'city_change_wt_discounting']=d['city_df_brand'].loc[(d['city_df_brand']['gc_platform']==row['gc_platform']) & (d['city_df_brand']['bgr']==row['bgr']) & (d['city_df_brand']['gc_city']==row['gc_city']) ,'change_wt_discounting'].values[0]
    dis_tab=dis_tab[dis_tab['change_wt_discounting']<0]
    if dis_tab.empty:
        print('No discounting drop found')
        return pd.DataFrame(columns=['gc_platform', 'bgr', 'What happened?', 'Why it matters?',
       'wt_discounting_dr1', 'change_wt_discounting', 'share_dr1',
       'pct_change'])

    for index, row in dis_tab.iterrows():
        # Get unique brands for comparison
        comps = ppu_cat_comp.loc[
            (ppu_cat_comp['gc_platform'] == row['gc_platform']) & 
            (ppu_cat_comp['bgr'] == row['bgr']) & 
            (ppu_cat_comp['ppu_category'] == row['ppu_category']), 
            'brand'
        ].unique()
        comp_ms = None
        a_discount = None
        c_discount = None
        comp_share = None
        comp_ms_delta = None
        
        if len(comps) >= 1:
            comp_df = d['merged_comp'].loc[
                (d['merged_comp']['gc_platform'] == row['gc_platform']) & 
                (d['merged_comp']['bgr'] == row['bgr']) &  (d['merged_comp']['gc_city']==row['gc_city']) &
                (d['merged_comp']['brand'].isin(comps)) & 
                (d['merged_comp']['change_wt_discounting'] > 0) & (d['merged_comp']['pct_change']>0)
            ].sort_values(by='share_dr1', ascending=False).head(4).sort_values(by='change_wt_discounting', ascending=False).head(1)
            if comp_df.empty:
                comps = ppu_cat_comp.loc[
            (ppu_cat_comp['gc_platform'] == row['gc_platform']) & 
            (ppu_cat_comp['bgr'] == row['bgr']), 
            'brand'
        ].unique()
                comp_df = d['merged_comp'].loc[
                (d['merged_comp']['gc_platform'] == row['gc_platform']) & 
                (d['merged_comp']['bgr'] == row['bgr']) & (d['merged_comp']['gc_city']==row['gc_city']) &
                (d['merged_comp']['brand'].isin(comps)) & 
                (d['merged_comp']['change_wt_discounting'] > 0) & (d['merged_comp']['pct_change']>0)
            ].sort_values(by='share_dr1', ascending=False).head(4).sort_values(by='change_wt_discounting', ascending=False).head(1)
                if comp_df.empty:
                    comp_ms=None
                    c_discount=None
                    a_discount=None
                    comp_share=None
                    comp_ms_delta=None
                else:
                    comp_ms=comp_df['brand'].values[0]
                    c_discount=comp_df['change_wt_discounting'].values[0]
                    a_discount=comp_df['wt_discounting_dr1'].values[0]
                    comp_share=comp_df['share_dr1'].values[0]
                    comp_ms_delta=comp_df['pct_change'].values[0]

                
            else:
                comp_ms=comp_df['brand'].values[0]
                c_discount=comp_df['change_wt_discounting'].values[0]
                a_discount=comp_df['wt_discounting_dr1'].values[0]
                comp_share=comp_df['share_dr1'].values[0]
                comp_ms_delta=comp_df['pct_change'].values[0]

        elif len(comps) == 0:
            continue  # Skip if no competitors found
        
        # Assigning comp_ms to the DataFrame
        dis_tab.loc[index, 'comp'] = comp_ms
        dis_tab.loc[index, 'comp_current_discount'] = a_discount
        dis_tab.loc[index, 'comp_delta_discount'] = c_discount
        dis_tab.loc[index, 'comp_share'] = comp_share
        dis_tab.loc[index, 'comp_ms_delta'] = comp_ms_delta


    
    rows = []
    for _, row in dis_tab.iterrows():
        # Extract own brand discount info
        our_current = row.get('wt_discounting_dr1', None)
        our_change = row.get('change_wt_discounting', None)
        
        # Extract competitor info
        comp = row.get('comp', '')
        comp_current = row.get('comp_current_discount', None)
        comp_change = row.get('comp_delta_discount', None)
        
        # Extract market share info
        our_ms = row.get('share_dr1', None)
        our_ms_change = row.get('pct_change', None)
        comp_share = row.get('comp_share', None)
        comp_ms_delta = row.get('comp_ms_delta', None)
        
        # Build the "What happened?" text
        what_happened = ""
        # if our_current is not None and our_change is not None:
        #     what_happened = f"Our discounting dropped to {our_current:.1f}% ({our_change:+.0f}%)"
        if comp and comp_current is not None and comp_change is not None:
            what_happened += f"\n{comp} increased discounting to {comp_current:.0f}% ({comp_change:+.0f}%) in {row['gc_city']}"
        
        # Build the "Why it matters?" text
        why_matters = ""
        if comp and comp_share is not None and comp_ms_delta is not None:
            why_matters = f"{comp}'s MS increased to {comp_share:.1f}% ({comp_ms_delta:+.0f}%) in {row['gc_city']}"
        # if our_ms is not None and our_ms_change is not None:
        #     why_matters += f"\nOur MS dropped to {our_ms:.1f}% ({our_ms_change:+.0f}%)"
        
        # Placeholder for detailed reference (fill as needed)
        detail_ref = ""
        
        rows.append({
            'gc_platform': row['gc_platform'],
            'bgr': row['bgr'],
            'What happened?': what_happened,
            'Why it matters?': why_matters,
            'wt_discounting_dr1':f"{row['cat_wt_discounting_dr1']:.1f}%",
            'change_wt_discounting':f"{row['cat_change_wt_discounting']:.1f}%",
            'share_dr1':f"{row['cat_share_dr1']:.1f}%",
            'pct_change':f"{row['cat_change_market_share']:.1f}%"
            # 'For details refer': detail_ref
        })
    
    return pd.DataFrame(rows)

def create_what_happened(d):
    """
    d: dict with keys
      - 'sku_own'     : SKU-level data for your own brand (DataFrame)
      - 'sku_comp'    : SKU-level data for competitor brands (DataFrame)
      - 'below2'      : brand-level metrics for your own brand (DataFrame)
      - 'filter_comp' : brand-level metrics for competitors (DataFrame)

    Returns a DataFrame with columns:
      ['platform', 'bgr_city', 'what_happened1', 'what_happened2', 'performance']

    For own-brand blocks, picks the top-2 metrics that decreased most; for competitor blocks, picks the top-2 metrics that increased most.
    Falls back to the top-2 by absolute change if not enough decreasing/increasing metrics exist.
    Adds a single 'performance' column where each row is all performance summaries joined by newlines.
    """
    # unpack and tag
    df_own_sku    = d['sku_own']
    df_comp_sku   = d['sku_comp']
    df_own_brand  = d['city_df_brand'].copy()
    df_comp_brand = d['filter_comp'].copy()
    df_own_brand['is_own_brand']  = True
    df_comp_brand['is_own_brand'] = False

    # label mappings
    brand_changes = {
        'pct_change':             'MS',
        'change_wt_osa':          'OSA',
        'change_wt_discounting':  'Discounting',
        'change_sov':             'overall SOV',
        'change_ogsov':           'Organic SOV',
        'change_adsov':           'Ad SOV'
    }
    sku_changes = {
        'change_market_share':    'MS',
        'change_wt_osa':          'OSA',
        'change_wt_discounting':  'Discounting',
        'change_sov':             'overall SOV',
        'change_ad_sov':          'Ad SOV',
        'change_organic_sov':     'Organic SOV'
    }
    # metric columns mapping
    brand_metric = {
        'pct_change':             'share_dr1',
        'change_wt_osa':          'wt_osa_dr1',
        'change_wt_discounting':  'wt_discounting_dr1',
        'change_sov':             'sov_dr1',
        'change_ogsov':           'organic_sov_dr1',
        'change_adsov':           'ad_sov_dr1'
    }
    sku_metric = {k: k.replace('change_','') + '_dr1' for k in sku_changes.keys()}

    rows = []
    for (plat, bgr, city,brand), own_grp in df_own_brand.groupby(['gc_platform','bgr','gc_city','brand']):
        key = {'gc_platform': plat, 'bgr':bgr,'gc_city':city,'brand':brand}
        comp_grp = df_comp_brand[
            (df_comp_brand['gc_platform']==plat) &
            (df_comp_brand['bgr']==bgr) &
            (df_comp_brand['gc_city']==city) &
            df_comp_brand['pct_change'].notna()
        ]
        blocks = pd.concat([own_grp, comp_grp], ignore_index=True)
        wh_texts = []
        perf_texts = []

        for _, br in blocks.iterrows():
            # select brand metrics
            others = {c: br.get(c, 0) for c in brand_changes if (c != 'pct_change') and (c!='change_ogsov')}
            if br['is_own_brand']:
                sel = sorted([(c, v) for c, v in others.items() if v < 0], key=lambda x: x[1])[:2]
                
                status = 'decreased'
                if len(sel)== 0:
                    fallback = [item for item in sorted(others.items(), key=lambda x: abs(x[1]), reverse=True) if item not in sel]
                    sel += fallback[:2-len(sel)]
                    status = 'increased'
            else:
                sel = sorted([(c, v) for c, v in others.items() if v > 0], key=lambda x: x[1], reverse=True)[:2]
                status = 'increased'
                if len(sel)== 0:
                    fallback = [item for item in sorted(others.items(), key=lambda x: abs(x[1]), reverse=True) if item not in sel]
                    sel += fallback[:2-len(sel)]
                    status = 'decreased'
            
            # performance summary
            labels = [brand_changes[c] for c, _ in sel]
            perf_texts.append(f"{br['brand']} " + ' and '.join(labels) + f" {status}")

            # brand summary text
            labels = ["MS"] + [brand_changes[c] for c,v in sel]

            # 2) compute the max length once
            upper_labels = [lbl.upper() for lbl in labels]

            max_label_len = max(len(lbl) for lbl in upper_labels)

            # add 1 or 2 extra spaces if you want padding between name & value
            name_width = max_label_len + 1
            mc = br.get('pct_change', 0)
            # own_val  = round(br.get('share_dr1', 0), 1)
            # own_dir  = '↑' if mc>0 else '↓'
            # own_str  = f"{own_val:.1f}%({own_dir}{abs(mc):.1f}%)"

            # sel_vals = [
            #     f"{round(br.get(brand_metric[c],0),1):.1f}%"
            #     f"({'↑' if v>0 else '↓'}{abs(v):.1f}%)"
            #     for c,v in sel
            # ]
            # val_width = max(len(own_str), *(len(s) for s in sel_vals)) + 1

            # # 4) redefine aligned_metric with those widths
            # def aligned_metric(name, val, change):
            #     direction = '↑' if change > 0 else '↓'
            #     val_str   = f"{val:.1f}%({direction}{abs(change):.1f}%)"
            #     return f"{name}\t{val_str}"

            # parts = [aligned_metric("MS", round(br.get('share_dr1', 0), 1), mc)]
            # for c, v in sel:
            #     val = round(br.get(brand_metric[c], 0), 1)
            #     parts.append(aligned_metric(brand_changes[c], val, v))

            # brand_line = f"{br['brand']}\n" + '\n'.join(parts)


            parts = [f"MS {round(br.get('share_dr1',0),1)}% ({'↑' if mc>0 else '↓'}{abs(mc):.1f}%)"]
            for c, v in sel:
                parts.append(f"{brand_changes[c]} {round(br.get(brand_metric[c],0),1)}% ({'↑' if v>0 else '↓'}{abs(v):.1f}%)")
            brand_line = f"{br['brand']} " + '\n\n'.join(parts)

            

            


            # SKU summary text
            skus = df_own_sku if br['is_own_brand'] else df_comp_sku
            sku_df = skus[
                (skus['gc_platform']==plat) &
                (skus['bgr']==bgr) &
                (skus['gc_city']==city) &
                (skus['brand']==br['brand'])
            ].copy()
            sku_df = sku_df[sku_df['change_market_share'].notna()]
            if sku_df.empty:
                wh_texts.append(np.nan)
            else:
                sku_df['abs_ms'] = sku_df['change_market_share'].abs()
                top = sku_df.sort_values('abs_ms', ascending=False).iloc[0]
                sk_parts = [
                    f"MS {round(top.get(sku_metric['change_market_share'],0),1)}%({'↑' if top['change_market_share']>0 else '↓'}{abs(top['change_market_share']):.1f}%)"
                ]
                others_s = {c: top.get(c, 0) for c in sku_changes if (c != 'change_market_share') and (c!='change_organic_sov')}
                if br['is_own_brand']:
                    sel_s = sorted([(c, v) for c, v in others_s.items() if v < 0], key=lambda x: x[1])[:2]
                else:
                    sel_s = sorted([(c, v) for c, v in others_s.items() if v > 0], key=lambda x: x[1], reverse=True)[:2]
                if len(sel_s) < 2:
                    fallback_s = [item for item in sorted(others_s.items(), key=lambda x: abs(x[1]), reverse=True) if item not in sel_s]
                    sel_s += fallback_s[:2-len(sel_s)]
                for c, v in sel_s:
                    sk_parts.append(f"{sku_changes[c]} {round(top.get(sku_metric[c],0),1)}%({'↑' if v>0 else '↓'}{abs(v):.1f}%)")
                sku_line = f"{top['title']}({round(top['grammage'],0):.0f}) " + ', '.join(sk_parts)
                # wh_texts.append(f"{brand_line}\n{sku_line}")
                wh_texts.append(f"{brand_line}")

        # attach what_happened fields
        for i, text in enumerate(wh_texts, start=1):
            key[f"what_happened{i}"] = text
        # single performance column joined with newlines
        key['performance'] = '\n'.join(perf_texts)
        rows.append(key)

    return pd.DataFrame(rows)


def summarize_impact_by_city(d):
    """
    Builds a DataFrame with impact and geographic segmentation text,
    applying special logic if overall share did not decline, and
    ensuring the number of cities referenced matches each 'top N' case.
    """
    df6 = pd.DataFrame()

    for _, r in d['below2'].iterrows():
        # filter top-3 cities for this platform/bgr
        df5 = d['full_city_df']
        df5 = df5[
            (df5['gc_platform'] == r['gc_platform']) &
            (df5['bgr'] == r['bgr'])
        ]
        if df5.empty:
            continue

        impact = r['sales_impact']
        drop_ms = r['pct_change']
        share_m1 = r['share_dr1']
        # 1) Impact text
        if drop_ms >= 0:
            impact_text = (
                f"While Pan-India share did not decline, localized impact "
                f"in this segment led to a sales loss of INR "
                f"{round(abs(impact) / 100000, 1)} lakhs"
            )
        else:
            impact_text = (
                f"Sales Opportunity Loss: INR {round(abs(impact) / 100000, 1)} lakhs\n\n"
                f"Category Share: {round(share_m1, 1)}% ({round((drop_ms), 1)}%)\n"
            )

        # 2) Geo text logic
        df_city = df5.sort_values('sales_impact', ascending=True).head(3).copy()
        df_city['pct_contrib'] = (df_city['sales_impact'] / impact * 100).round(1)
        cities = df_city['gc_city'].tolist()[:1]
        pcts = df_city['pct_contrib'].tolist()[:1]
        total_pct = sum(pcts)

        geo_text = ""  # default

        # inside your geo_text logic …

        # Top 1
        if len(cities) >= 1 and pcts[0] >= 50:
            pct0 = int(round(pcts[0]))
            geo_text = f"{cities[0]}({pct0}%)"

        # Top 2
        elif len(cities) >= 2 and pcts[0] + pcts[1] >= 50:
            pct0, pct1 = int(round(pcts[0])), int(round(pcts[1]))
            combined = pct0 + pct1
            geo_text = (
                f"{cities[0]}({pct0}%), {cities[1]}({pct1}%)"
            )

        # Top 3
        elif len(cities) >= 3 and pcts[0] + pcts[1] + pcts[2] >= 50:
            pct0, pct1, pct2 = map(lambda x: int(round(x)), pcts[:3])
            combined = pct0 + pct1 + pct2
            geo_text = (
                f"{cities[0]}({pct0}%), {cities[1]}({pct1}%)"
                # f"{cities[2]} ({pct2}%)"
            )

        # Fallback
        else:
            int_pcts = [int(round(x)) for x in pcts]
            parts = ", ".join(f"{city} ({pct}%)" for city, pct in zip(cities, int_pcts))
            total_pct = sum(int_pcts)
            geo_text = f"{parts}"


        # Build row
        df7 = pd.DataFrame({
            'gc_platform': [r['gc_platform']],
            'bgr': [r['bgr']],
            'brand':[r['brand']],
            # use top city for gc_city field
            'impact_text': [impact_text],
            'sales_impact':[r['sales_impact']],
            'share_dr1':[r['share_dr1']],
            'pct_change':[r['pct_change']],
            'geo_text':[geo_text]
        })

        df6 = pd.concat([df6, df7], ignore_index=True)

    return df6



def get_vars(d):
    ms = d['ms'].copy()
    month_map = {
        1:'January',2:'February',3:'March',4:'April',
        5:'May',6:'June',7:'July',8:'August',
        9:'September',10:'October',11:'November',12:'December'
    }
    ms['month_name'] = ms['month'].map(month_map)
    brand = ms['account'].iat[0]

    # first / last month names & GMV change
    min_m, max_m = ms['month'].min(), ms['month'].max()
    month1 = ms.loc[ms['month']==min_m, 'month_name'].iat[0]
    month2 = ms.loc[ms['month']==max_m, 'month_name'].iat[0]

    tot_min = ms.loc[ms['month']==min_m, 'actual_sum'].sum() / 1e7
    tot_max = ms.loc[ms['month']==max_m, 'actual_sum'].sum() / 1e7
    delta_tot = round(100*(tot_max - tot_min)/tot_min, 1)

    # build platform × month table
    pt = (
        ms
        .groupby(['gc_platform','month_name'])['actual_sum']
        .sum()
        .unstack(fill_value=0)
        .reset_index()
    )[[ 'gc_platform', month1, month2 ]]

    # compute shares
    den1, den2 = pt[month1].sum(), pt[month2].sum()
    pt[f'share_{month1}'] = round(100 * pt[month1]/den1, 1)
    pt[f'share_{month2}'] = round(100 * pt[month2]/den2, 1)
    pt['share_delta']       = pt[f'share_{month2}'] - pt[f'share_{month1}']

    # how many platforms do we have?
    no_of_platforms = pt['gc_platform'].nunique()

    # choose up to N platforms (e.g. top 3), or all if fewer
    N = min(3, no_of_platforms)
    top_platforms = (
        pt
        .assign(combined = pt[month1] + pt[month2])
        .nlargest(N, 'combined')
        ['gc_platform']
        .tolist()
    )

    # build your vars dict
    vars = {
        'BRAND': brand,
        'month1': month1,
        'month2': month2,
        'GMV1': f"₹ {tot_min:.1f} Cr",
        'GMV2': f"₹ {tot_max:.1f} Cr",
        'GMV3': f"{'↑' if delta_tot>0 else '↓'} {abs(delta_tot)}%",
    }
    # print(pd.DataFrame([vars]).to_markdown())

    # dynamically add Pmom, P{i}m1, P{i}m2, P{i}m3
    for i, p in enumerate(top_platforms, start=1):
        vars[f'Pmom{i}'] = p
        s1 = pt.loc[pt.gc_platform==p, f'share_{month1}'].iat[0]
        s2 = pt.loc[pt.gc_platform==p, f'share_{month2}'].iat[0]
        sd = pt.loc[pt.gc_platform==p, 'share_delta'].iat[0]
        vars[f'P{i}m1'] = f"{s1}%"
        vars[f'P{i}m2'] = f"{s2}%"
        vars[f'P{i}m3'] = f"{'↑' if sd>0 else '↓'} {round(abs(sd),1)}%"

    return pd.DataFrame([vars])



def get_gain_df(acc,l1_start,l1_end,l2_start,l2_end,d):
    gaindf=pd.DataFrame()
    bgr_data=get_bgr_data(acc,l1_start,l1_end,l2_start,l2_end)
    cnt_df=get_contribution(acc,l1_start,l1_end,l2_start,l2_end)
    cnt_df.columns=cnt_df.columns.str.lower()
    cnt_df['gc_platform']=cnt_df['gc_platform'].str.lower()
    cnt_df['bgr']=cnt_df['bgr'].str.lower()
    def format_city_share(df,sales_impact):
    # Sort by highest pct_change and select top cities if needed
        df_sorted = df.sort_values('sales_impact', ascending=False)
        cities = df_sorted['gc_city'].tolist()
        changes = ((df_sorted['sales_impact']/sales_impact)*100).round(1).tolist()
        # ms=df_sorted['share_dr1'].round(1).tolist()

        n = len(cities)

        if n == 0:
            return ""
        elif n == 1:
            return f"{cities[0]} ({changes[0]}%)"
        elif n == 2:
            return f"{cities[0]} ({changes[0]}%) and {cities[1]} ({changes[1]}%)"
        else:
            city_str = ", ".join(cities[:-1]) + f", and {cities[-1]}"
            change_str = ", ".join([f"{c}%" for c in changes[:-1]]) + f", and {changes[-1]}%"
            return f"{city_str} gained max. Share by {change_str}, respectively."
    own_bgr_data=bgr_data[bgr_data['is_own_brand']==1]
    if d['top2'].empty:
        return pd.DataFrame(columns=['gc_platform', 'bgr', 'Result', 'Offtake', 'geo_text', 'ms',
       'What worked for us?'])
    df_gain=d['top2'].sort_values(by='sales_impact',ascending=False).merge(own_bgr_data,on=['gc_platform','bgr','brand','is_own_brand'],how='left')
    for index,row in df_gain.iterrows():
        gaindf.loc[index,'gc_platform']=row['gc_platform']
        gaindf.loc[index,'bgr']=row['bgr']
        top_cities=d['full_city_df_top'].loc[(d['full_city_df_top']['gc_platform']==row['gc_platform']) & (d['full_city_df_top']['bgr']==row['bgr']) & (d['full_city_df_top']['brand']==row['brand'])].head(2)
        top_str=format_city_share(top_cities,row['sales_impact'])
        offtake_str = f"{row['sales_impact']/1e7:.1f} Cr" if row['sales_impact'] >= 1e7 else f"{row['sales_impact']/1e5:.1f} L"

        gaindf.loc[index, 'Result'] = (
            f"Offtake: {offtake_str} ({100 * (row['actual_sales_dr1'] - row['actual_sales_dr2']) / row['actual_sales_dr2']:.1f}%)\n"
            f"Category Share: {row['market_share_dr1']}% ({row['pct_change']:.1f}%)\n"
            f"{top_str}"
        )
        gaindf.loc[index,'Offtake']=offtake_str
        
        # gaindf.loc[index,'Category Share']=f"{row['market_share_dr1']}% ({row['pct_change']:.1f}%)"
        gaindf.loc[index,'geo_text']=top_str
        gaindf.loc[index,'ms']=f"{row['market_share_dr1']}% ({row['pct_change']:.1f}%)"


        # gaindf.loc[index,'Result']=f"Offtake: {row['actual_sales_dr1'].apply(lambda x: f'{x/1e7:.1f} Cr' if x >= 1e7 else f'{x/1e5:.1f} L')}({100*(row['actual_sales_dr1']-row['actual_sales_dr2'])/row['actual_sales_dr2']:.1f}%)\nCategory Share: {row['market_share_dr1']}({row['pct_change']:.1f}%)\n{top_str}"


    kwdf=get_relevant_keywords(acc,l1_start,l1_end,l2_start,l2_end)[0]
    kwdf['change_adsov']=kwdf['sov_dr1']-kwdf['sov_dr2']
    kwdf=kwdf[(kwdf['relevant_score']>1) & kwdf['change_adsov']>0]
    df_gain=d['city_df_brand_top'].sort_values(by='sales_impact',ascending=False)
    for index,row in df_gain.iterrows():
    
        rec_text=[]
        if row['change_wt_osa'] > 0:
            gain_skus = (
                d['a']
                .loc[
                    (d['a']['gc_platform'] == row['gc_platform']) &
                    (d['a']['bgr'] == row['bgr']) &
                    (d['a']['brand'] == row['brand']) &
                    (d['a']['gc_city'] == row['gc_city'])
                ]
                .merge(cnt_df, on=['gc_platform', 'bgr', 'gc_city', 'product_id', 'title'], how='left')
            )
            # if row['gc_platform']=='zepto':
            #     print(gain_skus.to_markdown())
            gain_skus = (
                gain_skus[gain_skus['change_ds_listing'] > 0]
                .sort_values(by='contribution_pct', ascending=False)['title']
                .head(2)
                .tolist()
            )
            
            # print(gain_skus)
            if len(gain_skus) == 0:
                pass  # or skip appending anything
            elif len(gain_skus) == 1:
                    
                # print(gain_skus)
                rec_text.append(
                    f"•{acc['db_name'].capitalize()}'s OSA increased to {row['wt_osa_dr1']}%(+{row['change_wt_osa']:.0f}%) "
                    f"by fixing DS listing for {gain_skus[0]} in {row['gc_city']}"
                )
            else:  # 2 or more
                   
                # print(gain_skus)
                rec_text.append(
                    f"• {acc['db_name'].capitalize()}'s OSA increased to {row['wt_osa_dr1']}%(+{row['change_wt_osa']:.0f}%) "
                    f"by fixing DS listing for {gain_skus[0]} and {gain_skus[1]} in {row['gc_city']}"
                )

       
        if row['change_adsov']>0:
            kwdf.gc_platform=kwdf.gc_platform.str.lower()
            kwdf.bgr=kwdf.bgr.str.lower()
            gain_kwd=kwdf.loc[(kwdf['gc_platform']==row['gc_platform']) & (kwdf['bgr']==row['bgr']) & (kwdf['gc_city']==row['gc_city'])].sort_values(by='final_score',ascending=False).head(10).sort_values(by='change_adsov',ascending=False).head(2)['keyword'].tolist()
            len_kd=len(gain_kwd)
            gain_kwd=", ".join(gain_kwd)
            rec_text.append(f"• SOV increased to {row['ad_sov_dr1']:.1f}(+{row['change_adsov']:.1f}%) by targeting relevant keywords - {gain_kwd}" if len_kd>0 else "")

        if row['change_wt_discounting']>0:
            rec_text.append(f"• Discounting increased to {row['wt_discounting_dr1']}(+{row['change_wt_discounting']:.1f}%) ,resulting in market share increase to {row['share_dr1']}({row['pct_change']:.1f}%)")
        gaindf.loc[((gaindf['bgr']==row['bgr']) & (gaindf['gc_platform']==row['gc_platform'])),'What worked for us?']='\n'.join(rec_text)
    return gaindf


def format_drainers(drainers: pd.DataFrame) -> pd.DataFrame:
    prefix_map = {'blinkit': 'BD', 'instamart': 'ID', 'zepto': 'ZD'}
    marker_positions = [1,2,3,4,6,7,8,9,10]

    # Build one dict per slide
    slides = []
    for plat, grp in drainers.groupby('platform'):
        grp = grp.reset_index(drop=True)
        prefix = prefix_map.get(plat, plat[:2].upper())
        for slide_idx in (1, 2):
            pfx = f"{prefix}{slide_idx}_"
            row = grp.iloc[slide_idx-1] if slide_idx-1 < len(grp) else None
            d = {}
            for m in marker_positions:
                key = f"{pfx}{m}"
                if row is not None:
                    if   m == 1:  d[key] = row['bgr']
                    elif m == 2:  d[key] = row['sales_impact']
                    elif m == 3:  d[key] = f"{row['share_dr1']}% ({row['pct_change']:.1f}%)"
                    elif m == 4:  d[key] = row['geo_text']
                    elif m == 6:  d[key] = row['city']
                    elif m == 7:  d[key] = row.get('what_happened1','') or ""
                    elif m == 8:  d[key] = row.get('what_happened2','') or ""
                    elif m == 9:  d[key] = row.get('what_happened3','') or ""
                    elif m == 10: d[key] = row['recommendation']
                else:
                    d[key] = ""
            slides.append(d)

    # Merge all slide-dicts into one, preferring non-empty values
    final = {}
    for s in slides:
        for k, v in s.items():
            # only overwrite if v is not None/empty/NA
            if v not in (None, "") and not pd.isna(v):
                final[k] = v
            elif k not in final:
                final[k] = ""

    # Enforce column ordering
    ordered_prefixes = []
    for platform in ['blinkit','instamart','zepto']:
        pref = prefix_map[platform]
        for slide_idx in (1, 2):
            ordered_prefixes.append(f"{pref}{slide_idx}_")
    cols = []
    for pfx in ordered_prefixes:
        for m in marker_positions:
            cols.append(f"{pfx}{m}")

    # Build single-row DataFrame in that order
    single_row = pd.DataFrame([{c: final.get(c, "") for c in cols}])
    return single_row




def format_gainers(gainers: pd.DataFrame) -> pd.DataFrame:
    """
    gainers: DataFrame with columns:
       ['platform', 'bgr', 'Offtake',    # numeric
        'ms',        # e.g. "5.5% (+0.5%)"
        'Top Impacted Cities',           # string
        'What worked for us?']           # string

    Returns a single‐row DataFrame with columns:
      BG1_1…BG1_5, BG2_1…BG2_5,
      ID1_1…ID2_5, ZG1_1…ZG2_5
    filled (or blank) for each of up to two top gainers per platform.
    """

    # 1) map platform→prefix
    prefix_map = {'blinkit': 'BG', 'instamart': 'IG', 'zepto': 'ZG'}
    markers    = [1,2,3,4,5]

    # 2) helper to format currency
   

    # 3) build one dict per “slide”
    slides = []
    for plat, grp in gainers.groupby('gc_platform'):
        grp = grp.reset_index(drop=True)
        pref = prefix_map.get(plat, plat[:2].upper())
        for idx in (1, 2):          # two slots per platform
            pfx = f"{pref}{idx}_"
            row = grp.iloc[idx-1] if idx-1 < len(grp) else None
            d = {}
            for m in markers:
                key = f"{pfx}{m}"
                if row is not None:
                    if   m == 1:  d[key] = row['bgr']
                    elif m == 2:  d[key] = row['Offtake']
                    elif m == 3:  d[key] = row['ms']
                    elif m == 4:  d[key] = row['geo_text']
                    elif m == 5:  d[key] = row['What worked for us?']
                else:
                    d[key] = ""
            slides.append(d)

    # 4) merge slides into one final dict
    final = {}
    for s in slides:
        for k, v in s.items():
            # overwrite only if non‐blank
            if pd.notna(v) and v != "":
                final[k] = v
            elif k not in final:
                final[k] = ""

    # 5) enforce ordering
    ordered_prefixes = []
    for plat in ['blinkit','instamart','zepto']:
        pr = prefix_map[plat]
        for idx in (1,2):
            ordered_prefixes.append(f"{pr}{idx}_")

    cols = []
    for pfx in ordered_prefixes:
        for m in markers:
            cols.append(f"{pfx}{m}")

    # 6) build and return the single‐row DataFrame
    single = pd.DataFrame([{c: final.get(c, "") for c in cols}])
    return single


def get_new_launches(acc,l1_start,l1_end,l2_start,l2_end,d):
    # acc['db_name']='sleepyowl'
    df1=pd.DataFrame()
    df=sfFetch(f"""
    select gc_platform,bgr,product_id,title,first_seen,is_own_brand from 
(select gc_platform,product_id,first_seen 
from wholetruth.silver.dim_platform_sku_master 
where product_id in (select distinct product_id from {acc['db_name']}.aggregate.qcom_product_dimensions) and gc_platform in ('Blinkit','Instamart','Zepto')) b
left join (select gc_platform,product_id ,title,is_own_brand,bgr 
from {acc['db_name']}.aggregate.qcom_product_dimensions)a 
using(gc_platform,product_id) 
where first_seen >= date_trunc('month',dateadd('month',-1,DATE '{l1_end}')) """)
    df.columns=df.columns.str.lower()
    df['gc_platform']=df['gc_platform'].str.lower()
    df['bgr']=df['bgr'].str.lower()
    df=df[df['is_own_brand']==0]
    # print(d['a'].head().to_markdown())
    # return df
    df=df.merge(d['a'][['gc_platform','bgr','gc_city','brand','product_id','title','is_own_brand','change_market_share','market_share_dr1','change_wt_osa','wt_osa_dr1']],on=['gc_platform','bgr','product_id','title','is_own_brand'],how='left')
    df=df[df['change_wt_osa'] > 0].sort_values(by=['gc_platform', 'market_share_dr1'], ascending =[True, False]).groupby('gc_platform', as_index=False).head(1)
    df['wt_osa_dr1']=df['wt_osa_dr1'].round(1)
    df['change_wt_osa']=df['change_wt_osa'].round(1)

    return df


    

def format_new_launches(new_df: pd.DataFrame) -> pd.DataFrame:
    """
    Formats up to two new launches into NL_1…NL_12 markers.
    If new_df is empty, returns a single-row DataFrame with all blank markers.
    """
    # Initialize blanks for 12 markers
    blank_markers = {f"NL_{i}": "" for i in range(1, 15)}
    
    # If no launches, return blank row immediately
    if new_df.empty:
        return pd.DataFrame([blank_markers])

    df = new_df.reset_index(drop=True)
    final = {}

    # Populate up to two launches
    for i in range(2):
        offset = i * 7
        if i < len(df):
            row = df.iloc[i]
            final[f"NL_{offset+1}"] = row['gc_platform']
            final[f"NL_{offset+2}"] = row['bgr']
            final[f"NL_{offset+3}"] = row['gc_city']
            final[f"NL_{offset+4}"] = row['brand']
            final[f"NL_{offset+5}"] = row['title']
            final[f"NL_{offset+6}"] = f"{round(row['wt_osa_dr1'],1)}% ({row['change_wt_osa']:.1f}%)"
            final[f"NL_{offset+7}"] = f"{row['market_share_dr1']:.1f}%"
        else:
            # fill remaining slots with blanks
            for j in range(1, 8):
                final[f"NL_{offset+j}"] = ""

    # Ensure all markers 1–12 are present
    all_markers = {**blank_markers, **final}
    cols = [f"NL_{i}" for i in range(1, 15)]
    return pd.DataFrame([{c: all_markers[c] for c in cols}])


def get_drainers(final):
    drainers=pd.DataFrame()
    import re

    def first_match(lines, pattern, flags=re.IGNORECASE):
        """Return first regex group(1) from lines or None."""
        rx = re.compile(pattern, flags)
        for ln in lines:
            m = rx.search(ln)
            if m:
                return m.group(1).strip()
        return None

    for idx, row in final.iterrows():
        # 1) Copy core fields
        drainers.loc[idx, 'platform']       = row['gc_platform']
        drainers.loc[idx, 'bgr']            = row['bgr']
        drainers.loc[idx, 'city']           = row['gc_city']
        drainers.loc[idx, 'geo_text']       = row['geo_text']
        drainers.loc[idx, 'what_happened1'] = row['what_happened1']
        drainers.loc[idx, 'what_happened2'] = row['what_happened2']
        drainers.loc[idx, 'what_happened3'] = row['what_happened3']
        drainers.loc[idx, 'Why it matters'] = row['impact_text']
        drainers.loc[idx, 'sales_impact']   = f"{abs(row['sales_impact'])/100000:.1f} L"
        drainers.loc[idx, 'share_dr1']      = row['share_dr1']
        drainers.loc[idx, 'pct_change']     = row['pct_change']

        # 2) Read raw rec blocks safely
        rec_raw = (row.get("Recommendation_sov") or "").strip()
        osa_raw = (row.get("Recommendation_osa") or "").strip()
        rec_lines = rec_raw.splitlines()
        osa_lines = osa_raw.splitlines()

        # 3) Extract SKU names
        #    - SOV bullet line: "• <SKU>: ..."
        #    - OSA header line: "- <SKU>:"
        sku_sov = first_match(rec_lines, r'^\s*•\s*(.+?)\s*:')
        sku_osa = first_match(osa_lines, r'^\s*-\s*(.+?)\s*:')

        # 4) Split SOV into "restore" vs "other" lines
        restore_lines = []
        other_lines = []
        for ln in rec_lines:
            if re.search(r'restore\s+osa', ln, flags=re.IGNORECASE):
                restore_lines.append(ln.strip())
            elif ln.strip():
                other_lines.append(ln.strip())

        # 5) If there are restore lines, move them into OSA block
        if restore_lines:
            for rln in restore_lines:
                # Extract the "right side" after first colon, if present
                # rln examples: "• SKU: Restore OSA to 95% ..."  OR  "• Restore OSA ..."
                parts = rln.split(':', 1)
                if len(parts) == 2:
                    right = parts[1].strip()
                else:
                    # no colon → treat whole line as right
                    right = re.sub(r'^\s*[•-]\s*', '', rln).strip()

                if sku_sov and sku_osa and (sku_sov == sku_osa):
                    # SAME SKU → keep bullet without SKU prefix
                    # "• Restore OSA to XX% ..."
                    moved = f"• {right}" if not right.lower().startswith('restore osa') else f"• {right}"
                    osa_lines.append(moved)
                else:
                    # DIFFERENT SKU → add a new OSA SKU header with this SKU
                    # "- <sku_sov>: Restore OSA to XX% ..."
                    # Fallback if sku_sov is missing → keep bullet form
                    if sku_sov:
                        moved = f"- {sku_sov}: {right}"
                    else:
                        moved = f"• {right}"
                    osa_lines.append(moved)

        # 6) Rebuild cleaned SOV block (without the restore lines)
        rec_clean_lines = other_lines[:]

        # 7) Housekeeping: drop orphaned headers if needed
        def drop_orphan_header(lines, header_regex):
            # remove header if it's the only non-empty line
            nonempty = [l for l in lines if l.strip()]
            if len(nonempty) == 1 and re.match(header_regex, nonempty[0], flags=re.IGNORECASE):
                return []
            # also if header present but no other bullets below it, drop it
            if nonempty and re.match(header_regex, nonempty[0], flags=re.IGNORECASE):
                if len(nonempty) == 1:
                    return []
            return lines

        rec_clean_lines = drop_orphan_header(rec_clean_lines, r'^\s*sov\s+recommendation:?\s*$')
        osa_lines       = drop_orphan_header(osa_lines, r'^\s*osa\s+recommendations?:?\s*$')

        # 8) Tidy spacing: collapse >2 blank lines, trim
        def tidy(lines):
            txt = "\n".join(lines)
            txt = re.sub(r'\n{3,}', '\n\n', txt).strip()
            return txt

        osa_text_clean = tidy(osa_lines)
        rec_text_clean = tidy(rec_clean_lines)

        # 9) Write back
        drainers.loc[idx, "Recommendation_osa"] = osa_text_clean if osa_text_clean else None
        drainers.loc[idx, "Recommendation_sov"] = rec_text_clean if rec_text_clean else None

        # 10) Combined final recommendation (optional)
        combined = "\n\n".join([x for x in [osa_text_clean, rec_text_clean] if x])
        drainers.loc[idx, "recommendation"] = combined if combined else None

        
    return drainers[~((drainers['recommendation']=='')|(drainers['recommendation'].isna())|(drainers['recommendation']=='\n'))]


# streamlit_app.py
# -------------------------------------------------------------
# MTD Framework: select db(s), domain(s), date ranges; generate PPT/PDF
# -------------------------------------------------------------
# Assumptions:
# - Your environment already exposes these project helpers:
#   fetch_query_results, get_full_data, fetch_offtake_uploaded_till,
#   get_sku_osa_data, create_what_happened, get_summary, get_osa_summary,
#   get_sov_summary, summarize_discount_table, get_gain_df, summarize_impact_by_city,
#   get_new_launches, format_sov, format_osa, format_disc, format_gainers,
#   format_new_launches, get_drainers, format_drainers
# - The placeholder replacement + PDF conversion helpers are defined below.
#   (Works with libreoffice/soffice on PATH to export PDF; otherwise falls back.)
#
# How to run:
#   streamlit run streamlit_app.py



# ==========================================
# Placeholder replacement & PDF conversion
# ==========================================
PLACEHOLDER_RE = re.compile(r"\{\{([^}]+)\}\}")


def copy_run_format(src_run, dst_run):
    s, d = src_run.font, dst_run.font
    d.name = s.name
    d.size = s.size
    d.bold = s.bold
    d.italic = s.italic
    d.underline = s.underline
    if s.color is not None and s.color.type is not None:
        try:
            d.color.rgb = s.color.rgb
        except Exception:
            pass


def _iter_textframes(shape):
    if getattr(shape, "has_text_frame", False):
        yield shape.text_frame
    if getattr(shape, "has_table", False):
        for r in shape.table.rows:
            for c in r.cells:
                if c.text_frame:
                    yield c.text_frame
    # group = 6
    if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
        for shp in shape.shapes:
            yield from _iter_textframes(shp)


def _replace_in_paragraph_preserve_runs(paragraph, mapping: dict, leave_missing=True):
    runs = list(paragraph.runs)
    pieces = [r.text or "" for r in runs]
    combined = "".join(pieces)
    if not combined or "{{" not in combined:
        return

    # map char index -> source run index
    idx_map = []
    for i, t in enumerate(pieces):
        idx_map.extend([i] * len(t))

    segments, last = [], 0

    def add_seg(text, src_idx):
        if text:
            segments.append((text, src_idx))

    for m in PLACEHOLDER_RE.finditer(combined):
        s, e = m.span()
        if s > last:
            src_idx = idx_map[last] if last < len(idx_map) else (idx_map[-1] if idx_map else 0)
            add_seg(combined[last:s], src_idx)

        key = m.group(1).strip()
        repl = str(mapping.get(key, "" if not leave_missing else m.group(0)))
        if s < len(idx_map):
            src_idx = idx_map[s]
        elif s > 0 and (s - 1) < len(idx_map):
            src_idx = idx_map[s - 1]
        else:
            src_idx = 0
        add_seg(repl, src_idx)
        last = e

    if last < len(combined):
        src_idx = idx_map[last] if last < len(idx_map) else (idx_map[-1] if idx_map else 0)
        add_seg(combined[last:], src_idx)

    # rebuild runs
    for r in runs[::-1]:
        r.text = ""
        paragraph._p.remove(r._r)
    for text, src_idx in segments:
        new_run = paragraph.add_run()
        new_run.text = text
        if runs:
            src_idx = max(0, min(src_idx, len(runs) - 1))
            copy_run_format(runs[src_idx], new_run)


def _replace_in_textframe_preserve(tf, mapping: dict, leave_missing=True):
    for p in tf.paragraphs:
        _replace_in_paragraph_preserve_runs(p, mapping, leave_missing=leave_missing)


def collect_slide_placeholders(prs: Presentation):
    by_slide = {}
    for slide in prs.slides:
        keys = set()
        for shape in slide.shapes:
            for tf in _iter_textframes(shape):
                txt = tf.text or ""
                for m in PLACEHOLDER_RE.finditer(txt):
                    keys.add(m.group(1).strip())
        by_slide[slide.slide_id] = keys
    return by_slide


def delete_slide(prs: Presentation, slide):
    slide_id = slide.slide_id
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        if sldId.id == slide_id:
            prs.part.drop_rel(sldId.rId)
            sldIdLst.remove(sldId)
            break


def delete_slides_where_all_placeholders_empty(prs: Presentation, slide_keys: dict, mapping: dict):
    for slide in list(prs.slides)[::-1]:
        keys = slide_keys.get(slide.slide_id, set())
        if not keys:
            continue
        if all((mapping.get(k) is None) or (str(mapping.get(k)) == "") for k in keys):
            delete_slide(prs, slide)


def replace_placeholders_in_pptx(pptx_in: str, pptx_out: str, mapping: dict, leave_missing=False):
    prs = Presentation(pptx_in)
    per_slide_keys = collect_slide_placeholders(prs)
    for slide in prs.slides:
        for shape in slide.shapes:
            for tf in _iter_textframes(shape):
                _replace_in_textframe_preserve(tf, mapping, leave_missing=leave_missing)
    delete_slides_where_all_placeholders_empty(prs, per_slide_keys, mapping)
    prs.save(pptx_out)


def pptx_to_pdf(pptx_bytes: bytes) -> bytes:
    """Convert PPTX bytes -> PDF bytes using cloud API service (Zamzar).
    Returns PDF bytes; raises RuntimeError if conversion is not possible.
    
    Falls back to LibreOffice if available locally, otherwise uses Zamzar API.
    """
    # Try LibreOffice first if available (for local development)
    exe = shutil.which("soffice") or shutil.which("libreoffice")
    if exe:
        try:
            with tempfile.TemporaryDirectory() as td:
                td_path = Path(td)
                in_path = td_path / "in.pptx"
                out_dir = td_path / "out"
                out_dir.mkdir(exist_ok=True)
                in_path.write_bytes(pptx_bytes)
                cmd = [
                    exe,
                    "--headless",
                    "--convert-to",
                    "pdf:impress_pdf_Export",
                    "--outdir",
                    str(out_dir),
                    str(in_path),
                ]
                import subprocess
                run_res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.STDOUT, text=True)
                if run_res.returncode == 0:
                    pdf_path = out_dir / "in.pdf"
                    if not pdf_path.exists():
                        cands = sorted(out_dir.glob("*.pdf"), key=lambda p: p.stat().st_mtime, reverse=True)
                        if cands:
                            pdf_path = cands[0]
                        else:
                            raise RuntimeError("Expected PDF not produced by LibreOffice.")
                    return pdf_path.read_bytes()
        except Exception as e:
            print(f"LibreOffice conversion failed: {e}, trying cloud API...")
    
    # Fall back to cloud API (Zamzar)
    try:
        import streamlit as st
        try:
            api_key = st.secrets.get("zamzar_api_key", None)
        except (AttributeError, FileNotFoundError, KeyError):
            api_key = None
        
        if not api_key:
            api_key = os.getenv("ZAMZAR_API_KEY")
        
        if not api_key:
            raise RuntimeError(
                "Zamzar API key not found. Please set 'zamzar_api_key' in Streamlit Cloud secrets "
                "or ZAMZAR_API_KEY environment variable. "
                "Get a free API key at https://developers.zamzar.com/"
            )
        
        # Upload file to Zamzar
        import requests
        import time
        
        files = {'source_file': ('presentation.pptx', pptx_bytes, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
        data = {'target_format': 'pdf'}
        headers = {'Authorization': f'Basic {api_key}'}
        
        # Submit conversion job
        response = requests.post('https://api.zamzar.com/v1/jobs', files=files, data=data, headers=headers, auth=(api_key, ''))
        response.raise_for_status()
        job_id = response.json()['id']
        
        # Poll for completion
        max_attempts = 30
        for attempt in range(max_attempts):
            time.sleep(2)
            status_response = requests.get(f'https://api.zamzar.com/v1/jobs/{job_id}', auth=(api_key, ''))
            status_response.raise_for_status()
            job_data = status_response.json()
            
            if job_data['status'] == 'successful':
                file_id = job_data['target_files'][0]['id']
                # Download converted file
                download_response = requests.get(f'https://api.zamzar.com/v1/files/{file_id}/content', auth=(api_key, ''), stream=True)
                download_response.raise_for_status()
                return download_response.content
            elif job_data['status'] == 'failed':
                raise RuntimeError(f"Zamzar conversion failed: {job_data.get('failure', {}).get('message', 'Unknown error')}")
        
        raise RuntimeError("Zamzar conversion timed out")
        
    except ImportError:
        raise RuntimeError("requests library required for cloud PDF conversion. Install with: pip install requests")
    except Exception as e:
        raise RuntimeError(f"PDF conversion failed: {str(e)}. Please ensure ZAMZAR_API_KEY is set in Streamlit secrets or environment variables.")


# ==========================================
# Data helpers (wrappers around your project functions)
# ==========================================



def build_insights_df(acc_row: pd.Series,
                       l1_start: dt.date, l1_end: dt.date,
                       l2_start: dt.date, l2_end: dt.date) -> pd.DataFrame:
    """Replicates your current pipeline but returns the final wide dataframe (fulldf)."""
    acc = acc_row
    d = get_full_data(acc)
    platform_update_date = fetch_offtake_uploaded_till(acc)

    # SKU OSA data (DR1 only)
    sku_data = get_sku_osa_data(acc, l1_start, l1_end, l2_start, l2_end, platform_update_date)
    sku_data = sku_data[sku_data["date_period"] == "dr1"]

    # PSL per-store merge (own brand only)
    def get_psl_mrp_for_stores(acc, l1_start, l1_end, l2_start, l2_end, sku_data):
        pids = tuple(
            sfFetch(
                f"""
                select distinct product_id
                from {acc['db_name']}.silver.sku_base_for_locality_metrics
                where is_own_brand=1
                """
            )["PRODUCT_ID"].unique()
        )
        store_data = sfFetch(
            f"""
            select gc_platform, product_id, store_id,
                   sum(potential_sales_loss_mrp) psl_mrp
            from wholetruth.aggregate.daily_store_sku_potential_sales_loss
            where gc_crawl_date between DATE '{l1_start}' and DATE '{l1_end}'
              and product_id in {pids}
              and store_id is not null
            group by 1,2,3
            """
        )
        store_data.columns = store_data.columns.str.lower()
        store_data["gc_platform"] = store_data["gc_platform"].str.lower()
        sku_data = sku_data.merge(store_data, on=["gc_platform", "product_id", "store_id"], how="left")
        sku_data["psl_mrp"] = sku_data["psl_mrp"].fillna(0)
        return sku_data

    f = get_psl_mrp_for_stores(acc, l1_start, l1_end, l2_start, l2_end, sku_data)
    f1=f.groupby(['date_period','gc_platform'])['psl_mrp'].sum().reset_index()

    sku_data=sku_data[(sku_data['gc_city']!='Others') & sku_data['psl_tags_pf']!='medium']
    sku_data1=sku_data.groupby(['gc_platform','bgr','gc_city','product_id','title']).agg(high_psl_stores=('locality','nunique'),listed_in=('total_listed','sum')).reset_index()

    # Summaries & recommendations
    try:
        e = create_what_happened(d)
    except Exception:
        e = pd.DataFrame(columns=["platform", "bgr", "gc_city", "brand", "what_happened1", "what_happened2", "performance"])

    sumdf = get_summary(d, acc, l1_start, l1_end, l2_start, l2_end,f1).fillna("")

    try:
        osa_rec = get_osa_summary(d, f, acc, l1_start, l1_end, l2_start, l2_end)
        final_osa = osa_rec[~((osa_rec["recommendation"] == "") | (osa_rec["recommendation"].isna()) | (osa_rec["recommendation"] == "\n"))].head(2)
    except Exception:
        osa_rec = pd.DataFrame(
            columns=[
                "gc_platform",
                "bgr",
                "top_geo",
                "share_dr1",
                "pct_change",
                "wt_osa_dr1",
                "change_wt_osa",
                "what_happened",
                "potential_sales_loss",
                "recommendation",
            ]
        )
        final_osa = osa_rec.copy()

    sov_rec = get_sov_summary(d, acc, l1_start, l1_end, l2_start, l2_end)
    final_sov = sov_rec[~((sov_rec["Recommendation"] == "") | (sov_rec["Recommendation"].isna()) | (sov_rec["Recommendation"] == "\n"))].head(2)

    dis_rec = summarize_discount_table(acc, l1_start, l1_end, l2_start, l2_end, d)
    final_dis_rec = dis_rec.head(2)

    final_gain = get_gain_df(acc, l1_start, l1_end, l2_start, l2_end, d)
    final_gain = final_gain[(~(final_gain["What worked for us?"] == "") | (final_gain["What worked for us?"] == "\n"))]

    summary = summarize_impact_by_city(d)
    e = e.merge(summary, on=["gc_platform", "bgr", "brand"], how="left")
    final_df = e.copy().fillna("")
    final_df = final_df.sort_index(axis=1)
    for i in range(1, 4):
        if f"what_happened{i}" not in final_df.columns:
            final_df[f"what_happened{i}"] = ""

    final = (
        final_df
        .merge(sov_rec[["gc_platform", "bgr", "Recommendation"]], on=["gc_platform", "bgr"], how="left")
        .rename(columns={"Recommendation": "Recommendation_sov"})
        .merge(osa_rec[["gc_platform", "bgr", "recommendation"]], on=["gc_platform", "bgr"], how="left", suffixes=("", "_osa"))
        .rename(columns={"recommendation": "Recommendation_osa"})
    )

    new_launches = get_new_launches(acc, l1_start, l1_end, l2_start, l2_end, d)

    final.fillna("", inplace=True)
    sovdf = format_sov(final_sov)
    osadf = format_osa(final_osa)
    disdf = format_disc(final_dis_rec)
    gaindf = format_gainers(final_gain)
    newdf = format_new_launches(new_launches)

    drainers = get_drainers(final)
    drainers_df = format_drainers(drainers)

    fulldf = pd.concat([sumdf, osadf, sovdf, disdf, drainers_df, gaindf, newdf], axis=1)
    fulldf["Brand"] = acc.db_name.capitalize() if hasattr(acc, "db_name") else str(acc["db_name"]).capitalize()
    fulldf['m2']=l1_start.strftime('%b')
    fulldf['m1']=l2_start.strftime('%b')
    return fulldf

def df_to_mapping(df: pd.DataFrame) -> Dict[str, str]:
    """Map column -> string value from a single-row dataframe for placeholder fill."""
    if df is None or df.empty:
        return {}
    row = df.iloc[0].to_dict()
    return {str(k): ("" if pd.isna(v) else str(v)) for k, v in row.items()}

@st.cache_data(show_spinner=True, ttl=3600)  # Cache for 1 hour, show spinner
def load_config_df() -> pd.DataFrame:
    """Load configuration from database. Returns empty DataFrame on failure."""
    # Clear previous error
    if 'config_load_error' in st.session_state:
        del st.session_state.config_load_error
    if 'config_load_steps' in st.session_state:
        del st.session_state.config_load_steps
    
    steps = []
    try:
        # Check if database URL is configured (check both env vars and Streamlit secrets)
        from queryhelper import get_env_var
        db_url = get_env_var('fetch_query_url') or os.getenv('fetch_query_url')
        steps.append(f"Step 1: Checking fetch_query_url... {'✅ Found' if db_url else '❌ Not found'}")
        if not db_url:
            st.session_state.config_load_steps = steps
            return pd.DataFrame()
        
        # Show connection URL (masked for security)
        from urllib.parse import urlparse
        try:
            parsed = urlparse(db_url)
            # Mask sensitive parts but show database name
            db_name_from_url = parsed.path.split('/')[-1] if parsed.path else 'unknown'
            masked_url = f"{parsed.scheme}://{parsed.hostname}:{parsed.port or 'default'}/{db_name_from_url}"
            steps.append(f"Step 2: Connection URL: `{masked_url}` (masked)")
            steps.append(f"   Database name from URL: `{db_name_from_url}`")
        except Exception as parse_error:
            steps.append(f"Step 2: Connection URL configured (cannot parse: {str(parse_error)})")
        
        # Attempting to connect to database - first get current schema
        steps.append("Step 3: Getting database and schema info...")
        try:
            db_info = fetch_query_results("SELECT current_database() as db_name, current_schema() as schema_name, current_schemas(false) as search_path")
            if not db_info.empty:
                db_name = db_info.iloc[0].get('db_name', 'Unknown')
                schema_name = db_info.iloc[0].get('schema_name', 'Unknown')
                steps.append(f"   Connected to database: `{db_name}`, current schema: `{schema_name}`")
        except:
            steps.append("   Could not get database info")
        
        # Try to find accounts table in different schemas
        steps.append("Step 4: Searching for accounts table...")
        found_schemas = []
        try:
            table_search = fetch_query_results("""
                SELECT table_schema, table_name
                FROM information_schema.tables
                WHERE table_name = 'accounts'
                AND table_schema NOT IN ('information_schema', 'pg_catalog', 'pg_toast')
                ORDER BY table_schema
            """)
            if not table_search.empty:
                found_schemas = table_search['table_schema'].unique().tolist()
                steps.append(f"   ✅ Found `accounts` table in schema(s): {', '.join(found_schemas)}")
            else:
                steps.append("   ⚠️ `accounts` table not found in information_schema")
                # Try to list all available schemas
                try:
                    all_schemas = fetch_query_results("""
                        SELECT schema_name 
                        FROM information_schema.schemata
                        WHERE schema_name NOT IN ('information_schema', 'pg_catalog', 'pg_toast', 'pg_temp_1', 'pg_toast_temp_1')
                        ORDER BY schema_name
                    """)
                    if not all_schemas.empty:
                        schema_list = all_schemas['schema_name'].unique().tolist()
                        steps.append(f"   Available schemas: {', '.join(schema_list[:10])}{'...' if len(schema_list) > 10 else ''}")
                except:
                    pass
                
                # List some tables that DO exist to help debug
                try:
                    existing_tables = fetch_query_results("""
                        SELECT table_schema, table_name
                        FROM information_schema.tables
                        WHERE table_schema NOT IN ('information_schema', 'pg_catalog', 'pg_toast')
                        ORDER BY table_schema, table_name
                        LIMIT 20
                    """)
                    if not existing_tables.empty:
                        steps.append(f"   Sample tables found: {len(existing_tables)} tables (showing first 20)")
                        # Group by schema
                        for schema in existing_tables['table_schema'].unique()[:5]:
                            tables_in_schema = existing_tables[existing_tables['table_schema'] == schema]['table_name'].tolist()
                            steps.append(f"     - Schema `{schema}`: {', '.join(tables_in_schema[:5])}{'...' if len(tables_in_schema) > 5 else ''}")
                except:
                    pass
        except Exception as search_error:
            steps.append(f"   ⚠️ Could not search for table: {str(search_error)}")
        
        # Attempting to query accounts table - try multiple schemas
        steps.append("Step 5: Querying accounts table...")
        a = None
        schemas_to_try = found_schemas if found_schemas else ['public']  # Try found schemas first, then public
        
        for schema in schemas_to_try:
            try:
                schema_name = schema if schema else 'public'
                steps.append(f"   Trying schema: `{schema_name}`...")
                a = fetch_query_results(
                    f"""
                    SELECT id account_id, primary_email_domain
                    FROM {schema_name}.accounts
                    """
                )
                if a is not None and not a.empty:
                    steps.append(f"   ✅ Query succeeded with schema `{schema_name}`!")
                    break
                else:
                    steps.append(f"   ⚠️ Query executed but returned no rows from `{schema_name}.accounts`")
            except Exception as schema_error:
                steps.append(f"   ❌ Query failed for `{schema_name}.accounts`: {str(schema_error)[:100]}")
        
        # If still no results, try without schema (default search path)
        if a is None or a.empty:
            try:
                steps.append("   Trying without explicit schema (using search path)...")
                a = fetch_query_results(
                    """
                    SELECT id account_id, primary_email_domain
                    FROM accounts
                    """
                )
                if a is not None and not a.empty:
                    steps.append("   ✅ Query succeeded without explicit schema!")
            except Exception as default_error:
                steps.append(f"   ❌ Default query also failed: {str(default_error)[:100]}")
        
        # Check if None or empty
        if a is None or a.empty:
            steps.append("❌ Query returned no accounts (table may be empty, in different schema, or query failed)")
            steps.append("💡 **Possible solutions:**")
            steps.append("   - Check if you're connected to the correct database")
            steps.append("   - Verify the table exists and has data")
            steps.append("   - Try specifying schema explicitly (e.g., `schema.accounts`)")
            st.session_state.config_load_steps = steps
            return pd.DataFrame()
        
        steps.append(f"✅ Found {len(a)} account(s)")
        # Found accounts
        a["domain"] = a["primary_email_domain"].str.split(".").str[0]

        steps.append("Step 6: Querying configs table...")
        cfg = fetch_query_results(
            """
            SELECT workspace_id, account_id, config->>'model_name' AS db_name
            FROM configs
            WHERE config_type = 'onboarding'
            """
        )
        # Check if None or empty
        if cfg is None or cfg.empty:
            steps.append("❌ Query returned no configurations (table may be empty or query failed)")
            st.session_state.config_load_steps = steps
            return pd.DataFrame()
        
        steps.append(f"✅ Found {len(cfg)} configuration(s)")
        # Found configurations
        cfg = cfg.merge(a, on="account_id", how="left")
        cfg['build_blinkit']=True
        cfg['build_instamart']=True
        cfg['build_zepto']=True
        steps.append("Step 7: ✅ Successfully merged and loaded configuration")
        st.session_state.config_load_steps = steps
        # Successfully loaded configurations
        return cfg
    except Exception as e:
        # Configuration load error - will be shown in UI diagnostics
        steps.append(f"❌ Exception occurred: {str(e)}")
        st.session_state.config_load_error = str(e)
        st.session_state.config_load_steps = steps
        import traceback
        st.session_state.config_load_traceback = traceback.format_exc()
        return pd.DataFrame()


# ==========================================
# Streamlit UI
# ==========================================

st.title("MTD Framework")

# with st.sidebar:
#     st.subheader("Template & Output")
#     template_file = st.file_uploader("Upload PPTX template", type=["pptx"], accept_multiple_files=False)
#     output_kind = st.radio("Output format", ["PPTX", "PDF"], horizontal=True)
#     leave_missing = st.checkbox("Leave missing placeholders as-is", value=False)
with st.sidebar:
    st.subheader("Template & Output")
    # Use template from repository
    template_path = Path(__file__).parent / "Template" / "28 July Run_V2.pptx"
    if not template_path.exists():
        st.error(f"Template not found at: {template_path}")
        st.stop()
    output_kind = st.radio("Output format", ["PPTX", "PDF"], horizontal=True)
    leave_missing = st.checkbox("Leave missing placeholders as-is", value=False)


# Load configuration (cached, spinner shown by cache decorator)
cfg = load_config_df()

# Diagnostic section (only show if config failed to load)
if cfg.empty or "domain" not in cfg.columns:
    with st.expander("🔍 Connection Diagnostics", expanded=True):
        # Test database connection
        st.subheader("Database Connection Test")
        test_button = st.button("🔍 Test Database Connection", type="primary")
        
        if test_button:
            test_status = st.empty()
            test_status.info("Testing connection...")
            
            # Test fetch_query_url
            try:
                from queryhelper import fetch_query_results
                
                # Test 1: Basic connection
                test_result = fetch_query_results("SELECT 1 as test")
                if test_result.empty:
                    test_status.warning("⚠️ Connection succeeded but basic query returned no results")
                else:
                    test_status.success("✅ Basic connection test passed!")
                
                # Test 2: Check if accounts table exists (check all schemas)
                try:
                    # First, try to find the table in any schema
                    accounts_check = fetch_query_results("""
                        SELECT table_schema, table_name, 
                               (SELECT COUNT(*) FROM information_schema.columns 
                                WHERE table_schema = t.table_schema 
                                AND table_name = t.table_name) as column_count
                        FROM information_schema.tables t
                        WHERE table_name = 'accounts'
                        AND table_schema NOT IN ('information_schema', 'pg_catalog', 'pg_toast')
                        ORDER BY table_schema
                    """)
                    
                    if not accounts_check.empty:
                        # Table exists in at least one schema
                        schemas = accounts_check['table_schema'].unique().tolist()
                        schema_info = f"Found in schema(s): {', '.join(schemas)}"
                        
                        # Try to get row count (will use default schema or search path)
                        try:
                            row_count = fetch_query_results("SELECT COUNT(*) as count FROM accounts")
                            if not row_count.empty:
                                count = int(row_count.iloc[0]['count']) if 'count' in row_count.columns else 0
                                if count > 0:
                                    test_status.success(f"✅ `accounts` table exists ({schema_info}) with {count} row(s)")
                                else:
                                    test_status.warning(f"⚠️ `accounts` table exists ({schema_info}) but is empty")
                            else:
                                test_status.warning(f"⚠️ `accounts` table exists ({schema_info}) but could not get row count")
                        except Exception as count_error:
                            test_status.warning(f"⚠️ `accounts` table exists ({schema_info}) but query failed: {str(count_error)}")
                    else:
                        # Table not found in information_schema, but let's try querying it anyway
                        # (it might exist but not be visible in information_schema due to permissions)
                        test_status.info("ℹ️ `accounts` table not found in information_schema, but will try querying it directly...")
                        try:
                            row_count = fetch_query_results("SELECT COUNT(*) as count FROM accounts")
                            if not row_count.empty:
                                count = int(row_count.iloc[0]['count']) if 'count' in row_count.columns else 0
                                test_status.success(f"✅ `accounts` table exists (accessible via query) with {count} row(s)")
                            else:
                                test_status.warning("⚠️ Could query `accounts` table but got no row count")
                        except Exception as direct_query_error:
                            test_status.error(f"❌ `accounts` table does not exist or is not accessible: {str(direct_query_error)}")
                except Exception as table_error:
                    test_status.warning(f"⚠️ Could not check accounts table: {str(table_error)}")
                
                # Test 3: Get database info
                try:
                    db_info = fetch_query_results("SELECT current_database() as db_name, current_schema() as schema_name")
                    if not db_info.empty:
                        db_name = db_info.iloc[0]['db_name'] if 'db_name' in db_info.columns else 'Unknown'
                        schema_name = db_info.iloc[0]['schema_name'] if 'schema_name' in db_info.columns else 'Unknown'
                        test_status.info(f"📊 Connected to database: `{db_name}` (schema: `{schema_name}`)")
                except:
                    pass
                
                # Test 4: Check all tables in the database
                try:
                    all_tables = fetch_query_results("""
                        SELECT table_schema, table_name, 
                               (SELECT COUNT(*) FROM information_schema.columns 
                                WHERE table_schema = t.table_schema 
                                AND table_name = t.table_name) as column_count
                        FROM information_schema.tables t
                        WHERE table_schema NOT IN ('information_schema', 'pg_catalog', 'pg_toast')
                        ORDER BY table_schema, table_name
                        LIMIT 20
                    """)
                    if not all_tables.empty:
                        with st.expander("📋 View available tables in database"):
                            st.dataframe(all_tables)
                            st.caption(f"Showing first 20 tables. Total: {len(all_tables)}")
                except Exception as tables_error:
                    test_status.warning(f"⚠️ Could not list tables: {str(tables_error)}")
                
                # Test 5: Try the actual query with more details
                try:
                    # First check row count
                    count_query = fetch_query_results("SELECT COUNT(*) as total FROM accounts")
                    total_rows = 0
                    if not count_query.empty:
                        total_rows = int(count_query.iloc[0]['total']) if 'total' in count_query.columns else 0
                    
                    actual_query = fetch_query_results("""
                        SELECT id account_id, primary_email_domain
                        FROM accounts
                        LIMIT 5
                    """)
                    if not actual_query.empty:
                        test_status.success(f"✅ Query successful! Found {len(actual_query)} account(s) (Total in table: {total_rows})")
                        with st.expander("View sample data"):
                            st.dataframe(actual_query)
                    else:
                        if total_rows > 0:
                            test_status.warning(f"⚠️ Query executed but returned no rows (table has {total_rows} rows but query returned empty)")
                            st.info("💡 This might indicate a schema/column name mismatch. Check if column names are correct.")
                        else:
                            test_status.warning("⚠️ Query executed but returned no rows (table is empty)")
                            st.info("💡 The `accounts` table exists but has no data. This might be expected if you're connecting to a different database than your local one.")
                except Exception as query_error:
                    test_status.error(f"❌ Query failed: {str(query_error)}")
                    st.exception(query_error)
                    st.info("💡 **Possible issues:**\n"
                           "- Column names might be different (e.g., `id` vs `account_id`)\n"
                           "- Table might be in a different schema\n"
                           "- Connection might be to a different database than local")
                    
            except Exception as e:
                test_status.error(f"❌ `fetch_query_url` connection failed: {str(e)}")
                st.exception(e)
        
        # Show configuration load steps and errors
        if 'config_load_steps' in st.session_state:
            st.subheader("Configuration Load Steps")
            with st.expander("View detailed steps", expanded=True):
                for step in st.session_state.config_load_steps:
                    st.write(step)
                
                # Show database comparison if available
                if 'connected_db' in st.session_state:
                    st.info(f"💡 **Connected Database:** `{st.session_state.connected_db}` (Schema: `{st.session_state.connected_schema}`)\n\n"
                           f"Compare this with your local database name. If they're different, your connection URL might be pointing to a different database.")
        
        if 'config_load_error' in st.session_state:
            st.error(f"**Configuration Load Error:** {st.session_state.config_load_error}")
            if 'config_load_traceback' in st.session_state:
                with st.expander("View full error traceback"):
                    st.code(st.session_state.config_load_traceback)
        
        st.info("💡 **Troubleshooting:**\n"
                "1. Ensure `.env.encrypted` is in your repository\n"
                "2. Add `encryption_key` to Streamlit Cloud secrets (Settings → Secrets)\n"
                "3. Verify the encryption key matches the one used to encrypt the file\n"
                "4. Check that your `.env` file contains `fetch_query_url` and `cms_url` before encryption\n"
                "5. Click 'Test Database Connection' above to verify the connection works\n"
                "6. Ensure your database allows connections from Streamlit Cloud IPs\n"
                "7. **If query returns no rows but works locally:** The connection URL might point to a different database. Check the database name in the connection test above.\n"
                "8. **If configuration not loaded:** Check if the `accounts` and `configs` tables exist and have data in the database shown above.")

# Option pickers
col1, col2, col3 = st.columns([1.2, 1, 1])

with col1:
    # ---- Domain → db_name auto-select ----
    if "prev_domains" not in st.session_state:
        st.session_state.prev_domains = []
    if "sel_dbs" not in st.session_state:
        st.session_state.sel_dbs = []

    # Check if config is loaded successfully
    if cfg.empty or "domain" not in cfg.columns:
        # Show more helpful error message
        error_msg = "⚠️ Configuration not loaded. Please check your database connection."
        if 'config_load_error' in st.session_state:
            error_msg += f"\n\n**Error details:** {st.session_state.config_load_error}"
        st.error(error_msg)
        st.info("💡 Expand the 'Connection Diagnostics' section above for detailed troubleshooting.")
        
        # Don't stop - let user see diagnostics and test connection
        # st.stop()  # Commented out so diagnostics are visible
    
    # Only show domain selector if config is loaded
    if cfg.empty or "domain" not in cfg.columns:
        st.warning("⚠️ Cannot load domains - configuration not available. Please check the Connection Diagnostics section above.")
        all_domains = []
    else:
        all_domains = sorted([d for d in cfg["domain"].dropna().unique().tolist() if d])
    
    sel_domains = st.multiselect("Domain(s)", options=all_domains, default=[], key="sel_domains", disabled=len(all_domains) == 0)

    if cfg.empty or "domain" not in cfg.columns:
        db_opts = []
    elif sel_domains:
        db_opts = sorted(cfg[cfg["domain"].isin(sel_domains)]["db_name"].dropna().unique().tolist())
    else:
        db_opts = sorted(cfg["db_name"].dropna().unique().tolist())

    # Auto-select all dbs whenever domain selection changes
    domains_changed = sorted(st.session_state.prev_domains) != sorted(sel_domains)
    if domains_changed:
        st.session_state.sel_dbs = db_opts[:] if sel_domains else []
        st.session_state.prev_domains = sel_domains[:]

    # Keep selection within available options
    st.session_state.sel_dbs = [x for x in st.session_state.sel_dbs if x in db_opts]

    # db_name picker (auto-filled)
    sel_dbs = st.multiselect(
        "db_name(s)",
        options=db_opts,
        default=st.session_state.sel_dbs,
        key="sel_dbs"
    )


with col3:
    today = dt.date.today()
    # Two ranges: L1 (current window) and L2 (comparison window)
    
    l1_start = st.date_input("L1 start", value=today.replace(day=1))
    l1_end = st.date_input("L1 end", value=today)
    l2_start = st.date_input("L2 start", value=(l1_start - dt.timedelta(days=30)))
    l2_end = st.date_input("L2 end", value=(l1_start - dt.timedelta(days=1)))
    l1_start = dt.datetime.combine(l1_start, dt.time.min)
    l1_end   = dt.datetime.combine(l1_end, dt.time.min)
    l2_start = dt.datetime.combine(l2_start, dt.time.min)
    l2_end   = dt.datetime.combine(l2_end, dt.time.min)


# Validation
if l1_end < l1_start:
    st.error("L1 end must be on/after L1 start.")
if l2_end < l2_start:
    st.error("L2 end must be on/after L2 start.")

st.markdown("---")

run = st.button("Generate")

if run:
    if not sel_dbs:
        st.warning("Please select at least one db_name.")
        st.stop()

    # Template path is already set above, just verify it exists
    if not template_path.exists():
        st.error(f"Template not found at: {template_path}")
        st.error("Please ensure Template/28 July Run_V2.pptx exists in the repository.")
        st.stop()
    with tempfile.TemporaryDirectory() as td:
        td_path = Path(td)
        outputs: List[tuple] = []  # (filename, bytes)
        progress = st.progress(0.0, text="Starting…")
        status = st.empty()

        # Filter cfg by selected dbs
        work_cfg = cfg[cfg["db_name"].isin(sel_dbs)].copy()
        total = len(work_cfg)

        for i, (_, acc_row) in enumerate(work_cfg.iterrows(), start=1):
            db = acc_row["db_name"]
            status.write(f"🔄 Processing **{db}** ({i}/{total}) - This may take 2-3 minutes…")

            try:
                # Show progress for data building
                with st.spinner(f"Building insights for {db}..."):
                    fulldf = build_insights_df(acc_row, l1_start, l1_end, l2_start, l2_end)
                    mapping = df_to_mapping(fulldf)

                # Fill PPTX
                with st.spinner(f"Generating presentation for {db}..."):
                    filled_pptx_path = td_path / f"{db}_MTD_{l1_start}_{l1_end}.pptx"
                    replace_placeholders_in_pptx(str(template_path), str(filled_pptx_path), mapping, leave_missing=leave_missing)
                    pptx_bytes = filled_pptx_path.read_bytes()

                if output_kind == "PDF":
                    try:
                        with st.spinner(f"Converting {db} to PDF..."):
                            pdf_bytes = pptx_to_pdf(pptx_bytes)
                        outputs.append((f"{db}_MTD_{l1_start}_{l1_end}.pdf", pdf_bytes))
                        status.write(f"✅ Completed **{db}** - PDF ready")
                    except Exception as ex:
                        st.warning(f"⚠️ PDF export failed for {db} (will provide PPTX instead): {ex}")
                        outputs.append((f"{db}_MTD_{l1_start}_{l1_end}.pptx", pptx_bytes))
                        status.write(f"✅ Completed **{db}** - PPTX ready")
                else:
                    outputs.append((f"{db}_MTD_{l1_start}_{l1_end}.pptx", pptx_bytes))
                    status.write(f"✅ Completed **{db}** - PPTX ready")

            except Exception as e:
                st.error(f"❌ Error processing {db}: {e}")
                status.write(f"❌ Failed to process **{db}**")

            progress.progress(i / max(total, 1), text=f"Processed {i}/{total} databases")

        status.write("Preparing downloads…")

        # Single vs multiple download
        if len(outputs) == 1:
            fname, blob = outputs[0]
            st.download_button(
                label=f"Download {fname}",
                data=blob,
                file_name=fname,
                mime=("application/pdf" if fname.endswith(".pdf") else "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
            )
        elif len(outputs) > 1:
            # zip them
            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
                for fname, blob in outputs:
                    zf.writestr(fname, blob)
            zip_buf.seek(0)
            st.download_button(
                label=f"Download {len(outputs)} files (ZIP)",
                data=zip_buf,
                file_name=f"mtd_outputs_{l1_start}_{l1_end}.zip",
                mime="application/zip",
            )
        else:
            st.info("No files were produced.")

    st.success("Done.")







